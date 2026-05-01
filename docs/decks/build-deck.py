"""
Build the Agentic UI overview deck — `agentic-ui-overview.pptx`.

Audience: senior execs (sections 1, 8.a, 9), architects (2-7, 8.b),
and developers (4-6, 8.c, 10). Slides are colour-coded by track in
the speaker-notes header so a presenter can mix-and-match.

Produces a self-contained `.pptx` next to this script. No external
templates required — the master is built programmatically from the
brand tokens below so the deck stays consistent if the theme changes.
"""
from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

OUT = Path(__file__).parent / "agentic-ui-overview.pptx"

# ── Brand tokens (v3 — professional refresh) ────────────────────────────────
# Restrained palette in the spirit of strategy / consulting decks: deep
# navy primary, a single warm-amber accent, broad grayscale ramp. Colour
# is reserved for status (ok / warn / bad) and structural emphasis;
# decoration is grayscale only.
BRAND        = RGBColor(0x1E, 0x3A, 0x5F)   # deep navy
BRAND_DEEP   = RGBColor(0x12, 0x24, 0x3D)
BRAND_TINT   = RGBColor(0xEC, 0xF1, 0xF6)
ACCENT       = RGBColor(0xC4, 0x82, 0x21)   # warm amber — emphasis only
ACCENT_TINT  = RGBColor(0xFA, 0xF0, 0xDD)
SURFACE      = RGBColor(0xFF, 0xFF, 0xFF)
SURFACE_2    = RGBColor(0xF7, 0xF7, 0xF5)   # warm off-white
SURFACE_3    = RGBColor(0xEE, 0xEE, 0xEA)
TEXT         = RGBColor(0x14, 0x18, 0x21)
TEXT_2       = RGBColor(0x36, 0x3D, 0x49)
MUTED        = RGBColor(0x6B, 0x73, 0x80)
FAINT        = RGBColor(0xA2, 0xAA, 0xB6)
BORDER       = RGBColor(0xE3, 0xE5, 0xE9)
DIVIDER      = RGBColor(0xEC, 0xEE, 0xF2)
OK           = RGBColor(0x2F, 0x6F, 0x4F)   # forest green
WARN         = RGBColor(0xC4, 0x82, 0x21)   # amber (= ACCENT)
BAD          = RGBColor(0xA8, 0x33, 0x2A)   # muted red
INFO         = RGBColor(0x35, 0x6B, 0x88)   # slate blue

W, H = Inches(13.333), Inches(7.5)            # 16:9 widescreen


def add(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.5)
    s.shadow.inherit = False
    return s


def rounded(slide, x, y, w, h, fill, line=None, radius=0.05):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = radius
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    s.shadow.inherit = False
    return s


def text_box(slide, x, y, w, h, text, *,
             size=14, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
             font="Calibri", anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def bullets(slide, x, y, w, h, items, *,
            size=14, color=TEXT, gap=4, marker="•", marker_color=BRAND):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.level = 0
        # marker
        r0 = p.add_run(); r0.text = marker + "  "
        r0.font.name = "Calibri"; r0.font.size = Pt(size); r0.font.bold = True
        r0.font.color.rgb = marker_color
        # body
        if isinstance(item, tuple):
            head, tail = item
            r1 = p.add_run(); r1.text = head + " "
            r1.font.name = "Calibri"; r1.font.size = Pt(size); r1.font.bold = True
            r1.font.color.rgb = color
            r2 = p.add_run(); r2.text = tail
            r2.font.name = "Calibri"; r2.font.size = Pt(size); r2.font.color.rgb = color
        else:
            r = p.add_run(); r.text = item
            r.font.name = "Calibri"; r.font.size = Pt(size); r.font.color.rgb = color


def code_block(slide, x, y, w, h, lines):
    rect(slide, x, y, w, h, RGBColor(0x0F, 0x17, 0x2A))
    tb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1),
                                  w - Inches(0.3), h - Inches(0.2))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(0)
        r = p.add_run(); r.text = line
        r.font.name = "Menlo"; r.font.size = Pt(11)
        r.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)


def chip(slide, x, y, label, *, fg=BRAND_DEEP, bg=BRAND_TINT, size=10, w=None):
    if w is None:
        w = Inches(max(0.55, 0.10 * len(label) + 0.4))
    h = Inches(0.32)
    rounded(slide, x, y, w, h, bg, radius=0.5)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.name = "Calibri"; r.font.size = Pt(size); r.font.bold = True
    r.font.color.rgb = fg
    return w


def slide_chrome(slide, eyebrow, title, audience_tag=None):
    """Top chrome (v3) — eyebrow row, large title, thin divider rule.

    Removed the brand stripe; restraint matters more for senior
    audiences than visual identity. Audience tag now sits inline with
    the eyebrow so the title gets full width.
    """
    # Eyebrow row
    text_box(slide, Inches(0.7), Inches(0.5), Inches(8), Inches(0.3),
             eyebrow, size=10, bold=True, color=MUTED)
    if audience_tag:
        text_box(slide, W - Inches(3.7), Inches(0.5), Inches(3.0), Inches(0.3),
                 "FOR " + audience_tag.upper(), size=10, bold=True,
                 color=ACCENT, align=PP_ALIGN.RIGHT)
    # Title
    text_box(slide, Inches(0.7), Inches(0.85), W - Inches(1.4), Inches(0.7),
             title, size=30, bold=True, color=TEXT)
    # Divider rule under the title
    rect(slide, Inches(0.7), Inches(1.5), Inches(1.0), Inches(0.025), BRAND)


def slide_footer(slide, idx, total, section):
    """Footer (v3) — single thin rule, page number on the right, section
    label on the left, no full-width band. Less visual weight."""
    # Hairline rule
    rect(slide, Inches(0.7), H - Inches(0.45), W - Inches(1.4),
         Inches(0.012), BORDER)
    text_box(slide, Inches(0.7), H - Inches(0.4), Inches(6), Inches(0.28),
             section, size=9, color=MUTED)
    text_box(slide, Inches(0.7), H - Inches(0.4), W - Inches(1.4),
             Inches(0.28),
             "Maverick — Agentic UI for Angular",
             size=9, color=FAINT, align=PP_ALIGN.CENTER)
    text_box(slide, W - Inches(2.1), H - Inches(0.4), Inches(1.4),
             Inches(0.28),
             f"{idx:02d} / {total:02d}", size=9, color=MUTED,
             align=PP_ALIGN.RIGHT)


def add_speaker_notes(slide, *blocks):
    notes = slide.notes_slide.notes_text_frame
    notes.text = blocks[0] if blocks else ""
    for b in blocks[1:]:
        p = notes.add_paragraph(); p.text = b


# ── Diagram primitives ──────────────────────────────────────────────────────

def arrow(slide, x1, y1, x2, y2, *, color=BRAND, weight=1.5, dashed=False, head_w=6, head_l=8):
    """Connector with an optional arrowhead at (x2, y2).

    Coordinates are EMU. Pass `head_w=0` to draw a plain line (used for
    sequence-diagram lifelines where the dashed line has no terminator).
    """
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    ln = line.line._get_or_add_ln()
    if head_w > 0:
        etree.SubElement(ln, qn('a:tailEnd'),
                         {'type': 'triangle', 'w': 'med', 'len': 'med'})
    if dashed:
        etree.SubElement(ln, qn('a:prstDash'), {'val': 'dash'})
    return line


def chip_node(slide, x, y, w, h, label, *, fill=BRAND_TINT, fg=BRAND_DEEP,
              border=None, size=11, bold=True):
    """A labeled node — used as building block for diagrams."""
    shape = rounded(slide, x, y, w, h, fill, line=border, radius=0.08)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.name = "Calibri"; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = fg
    return shape


def labeled_box(slide, x, y, w, h, title, body, *,
                accent=BRAND, fill=SURFACE, title_size=12, body_size=10):
    """A diagram tile — accent stripe on top, title in colour, body below."""
    rounded(slide, x, y, w, h, fill, line=BORDER, radius=0.04)
    rect(slide, x, y, w, Inches(0.05), accent)
    text_box(slide, x + Inches(0.12), y + Inches(0.12), w - Inches(0.24),
             Inches(0.32), title, size=title_size, bold=True, color=accent)
    text_box(slide, x + Inches(0.12), y + Inches(0.45), w - Inches(0.24),
             h - Inches(0.55), body, size=body_size, color=TEXT_2)


def lane(slide, x, y, w, h, label, *, color=BRAND, lane_label_w=Inches(1.2)):
    """A horizontal swimlane — used in sequence diagrams."""
    rect(slide, x, y, lane_label_w, h, color)
    text_box(slide, x, y, lane_label_w, h, label,
             size=11, bold=True, color=SURFACE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rounded(slide, x + lane_label_w, y, w - lane_label_w, h, SURFACE,
            line=BORDER, radius=0.02)


def numbered_step(slide, x, y, w, h, n, title, body, *, accent=BRAND):
    """A numbered card — used in sequence/lifecycle diagrams."""
    rounded(slide, x, y, w, h, SURFACE, line=BORDER, radius=0.04)
    rect(slide, x, y, Inches(0.06), h, accent)
    # Number badge
    badge_size = Inches(0.42)
    rounded(slide, x + Inches(0.18), y + Inches(0.15), badge_size, badge_size,
            accent, radius=0.5)
    text_box(slide, x + Inches(0.18), y + Inches(0.15), badge_size, badge_size,
             str(n), size=14, bold=True, color=SURFACE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text_box(slide, x + Inches(0.72), y + Inches(0.16), w - Inches(0.85),
             Inches(0.32), title, size=13, bold=True, color=TEXT)
    text_box(slide, x + Inches(0.72), y + Inches(0.5), w - Inches(0.85),
             h - Inches(0.6), body, size=10, color=TEXT_2)


def event_pill(slide, x, y, w, label, *, color=BRAND):
    """Tiny pill for an event in a timeline."""
    h = Inches(0.32)
    rounded(slide, x, y, w, h, color, radius=0.5)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.name = "Menlo"; r.font.size = Pt(9); r.font.bold = True
    r.font.color.rgb = SURFACE


def big_number(slide, x, y, w, h, value, caption, *, color=BRAND, size=64):
    """Large numeral with caption — strategy-deck staple. Reserves
    the slide's top-of-mind real estate for one striking metric."""
    text_box(slide, x, y, w, Inches(h.emu / 914400 * 0.6), str(value),
             size=size, bold=True, color=color, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP)
    text_box(slide, x, y + Inches(size / 64),
             w, h - Inches(size / 64), caption,
             size=12, color=MUTED, align=PP_ALIGN.LEFT)


def stat_card(slide, x, y, w, h, value, caption, *, color=BRAND):
    """Compact stat card — bordered, big number above caption."""
    rounded(slide, x, y, w, h, SURFACE, line=BORDER, radius=0.04)
    text_box(slide, x + Inches(0.15), y + Inches(0.18),
             w - Inches(0.3), Inches(0.7),
             str(value), size=32, bold=True, color=color,
             align=PP_ALIGN.LEFT)
    text_box(slide, x + Inches(0.15), y + h - Inches(0.45),
             w - Inches(0.3), Inches(0.3),
             caption, size=10, color=MUTED)


def quote_block(slide, x, y, w, h, body, attribution=None, *, color=BRAND):
    """Pull-quote callout — vertical accent rule + italic body."""
    rect(slide, x, y, Inches(0.04), h, color)
    text_box(slide, x + Inches(0.3), y + Inches(0.05),
             w - Inches(0.4), h - Inches(0.5),
             body, size=16, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
    if attribution:
        text_box(slide, x + Inches(0.3), y + h - Inches(0.4),
                 w - Inches(0.4), Inches(0.3),
                 "— " + attribution, size=11, color=MUTED)


def key_facts(slide, x, y, w, h, rows):
    """A clean two-column key-facts strip — label / value pairs.
    Used instead of bullet lists for executive summaries."""
    rh = h.emu / max(1, len(rows))
    for i, (label, value) in enumerate(rows):
        ty = y + Inches(rh / 914400) * i
        if i > 0:
            rect(slide, x, ty, w, Inches(0.012), DIVIDER)
        text_box(slide, x + Inches(0.1), ty + Inches(0.15),
                 Inches(3.5), Inches(0.4),
                 label, size=10, bold=True, color=MUTED)
        text_box(slide, x + Inches(3.7), ty + Inches(0.15),
                 w - Inches(3.9), Inches(0.4),
                 value, size=12, color=TEXT)


def actor_lifeline(slide, x, y_top, y_bottom, label, *, color=BRAND):
    """Sequence-diagram actor head + dashed lifeline going down."""
    head_w = Inches(1.5); head_h = Inches(0.5)
    rounded(slide, x - head_w / 2, y_top, head_w, head_h, color, radius=0.2)
    text_box(slide, x - head_w / 2, y_top, head_w, head_h, label,
             size=10, bold=True, color=SURFACE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    arrow(slide, x, y_top + head_h, x, y_bottom, color=BORDER,
          weight=1, dashed=True, head_w=0, head_l=0)


# ── Slide builders ──────────────────────────────────────────────────────────

def slide_cover(prs):
    """Cover (v3) — restrained, navy-on-white, big serene type. Two
    horizontal rules anchor the composition. No gradients."""
    s = add(prs)
    # Subtle warm-white bg
    rect(s, Inches(0), Inches(0), W, H, SURFACE_2)
    # Top wordmark row
    text_box(s, Inches(0.7), Inches(0.6), Inches(8), Inches(0.32),
             "MAVERICK", size=11, bold=True, color=BRAND)
    text_box(s, W - Inches(4.5), Inches(0.6), Inches(3.8), Inches(0.32),
             "AGENTIC UI · OVERVIEW DECK",
             size=10, color=MUTED, align=PP_ALIGN.RIGHT)
    rect(s, Inches(0.7), Inches(1.0), W - Inches(1.4), Inches(0.018), BRAND)

    # Title block
    text_box(s, Inches(0.7), Inches(2.6), W - Inches(1.4), Inches(1.2),
             "Agentic UI for Angular",
             size=54, bold=True, color=TEXT)
    text_box(s, Inches(0.7), Inches(3.5), W - Inches(1.4), Inches(0.6),
             "AG-UI · Hashbrown · A2UI · MCP",
             size=32, color=BRAND)

    # Subtitle / value prop
    text_box(s, Inches(0.7), Inches(4.45), Inches(9), Inches(0.5),
             "One library. Three protocols. Microfrontend-native.",
             size=16, color=TEXT_2)

    # Audience footer block — restrained text rows, not pills
    y = Inches(6.1)
    rect(s, Inches(0.7), y, W - Inches(1.4), Inches(0.018), BORDER)
    text_box(s, Inches(0.7), y + Inches(0.15), Inches(3), Inches(0.3),
             "PREPARED FOR", size=9, bold=True, color=MUTED)
    text_box(s, Inches(0.7), y + Inches(0.4), W - Inches(1.4), Inches(0.4),
             "Senior executives  ·  Solution architects  ·  Developers",
             size=14, color=TEXT)

    # Date / version line
    text_box(s, Inches(0.7), H - Inches(0.6), W - Inches(1.4), Inches(0.3),
             "Q2 2026 release  ·  v1.0  ·  Confidential",
             size=10, color=FAINT)

    add_speaker_notes(s,
        "Cover. Restrained on purpose — strategy/consulting deck convention.",
        "30-minute path: cover, 1, 2, 3, 6, 11, 13, 16, 17, 19, 21, 22, 23, 25, 26-28, 29.",
        "60-minute path: include the AG-UI / Hashbrown deep dives and the registry tour.")
    return s


def slide_agenda(prs):
    s = add(prs)
    slide_chrome(s, "AGENDA", "What you'll leave with")
    items = [
        ("The shift.", "Why agentic UIs are the next interface paradigm — and the protocol problem nobody owns."),
        ("Three protocols.", "AG-UI, Hashbrown, A2UI — origin, model, fit, gaps."),
        ("One library.", "@maverick/agentic-ui — 13 registries, MFE-native, MCP-ready, telemetry-instrumented."),
        ("Working examples.", "demo-monolith, demo-shell + remotes, demo-multi-agent, demo-ediscovery — all in this repo."),
        ("Decision frameworks.", "Protocol matrix, when-to-use guide, persona-by-persona benefit map."),
        ("Roadmap & CTA.", "Where it's going through v1.0 — and where you can plug in."),
    ]
    bullets(s, Inches(0.7), Inches(1.7), W - Inches(1.4), Inches(5.4), items, size=16, gap=8)
    add_speaker_notes(s, "Agenda. Pace this for 30 or 60 minutes — sub-deck markers in speaker notes for each track.")
    return s


def slide_shift(prs):
    s = add(prs)
    slide_chrome(s, "1 · THE SHIFT", "From chatbots to agentic UIs", "Senior executives")

    # Three columns: yesterday, today, tomorrow
    col_w = (W - Inches(1.4) - Inches(0.4)) / 3
    titles = ["Yesterday — chatbots", "Today — agentic UIs", "Tomorrow — agent-native apps"]
    bodies = [
        ["Text-only conversation",
         "Apps wired to LLMs as a fancy autocomplete",
         "Hand-coded handlers per use case",
         "No structured affordances — every flow re-debated"],
        ["Tool-call protocols (AG-UI / Hashbrown / A2UI)",
         "Generative UI: LLM emits widgets, host renders",
         "Multi-agent orchestration with sticky routing",
         "MCP server for analyst workstations"],
        ["UI itself is plug-in registry data",
         "Federated remotes contribute capabilities mid-session",
         "Audit-grade observability across LLM ↔ tool ↔ UI",
         "Persona-scoped governance baked in"],
    ]
    colors = [MUTED, BRAND_DEEP, OK]
    for i in range(3):
        x = Inches(0.7) + i * (col_w + Inches(0.2))
        rounded(s, x, Inches(1.7), col_w, Inches(5.0), SURFACE,
                line=BORDER, radius=0.04)
        rect(s, x, Inches(1.7), col_w, Inches(0.06), colors[i])
        text_box(s, x + Inches(0.25), Inches(1.85), col_w - Inches(0.5),
                 Inches(0.4), titles[i], size=15, bold=True, color=colors[i])
        bullets(s, x + Inches(0.25), Inches(2.4), col_w - Inches(0.5),
                Inches(4.0), bodies[i], size=12, gap=6, marker="—",
                marker_color=colors[i])
    add_speaker_notes(s,
        "The exec framing. Yesterday was wrappers around an LLM. Today's protocols (AG-UI/Hashbrown/A2UI) make the agent + UI a real surface.",
        "Tomorrow: the UI itself becomes plugin data. We are building for the third column.")
    return s


def slide_why_now(prs):
    s = add(prs)
    slide_chrome(s, "1 · THE SHIFT", "Why now", "Senior executives")

    drivers = [
        ("LLM-native tool calling matured.", "OpenAI, Anthropic, Google, Mistral all ship structured tool calls. Quality and latency now meet enterprise bars."),
        ("Three competing protocols emerged.", "AG-UI from CopilotKit, Hashbrown from Liquid Frontiers, A2UI from a multi-vendor consortium. Each is the leader in a slice."),
        ("MCP standardised the server side.", "Anthropic's Model Context Protocol made tool definitions portable across analyst workstations (Claude Desktop, Cursor, Zed)."),
        ("Generative UI moved past PoC.", "LLMs reliably emit JSON specifying widget + props; Angular signals + dynamic component loader make rendering trivial."),
        ("Regulated industries are catching up.", "Legal, healthcare, finance — all need defensible audit trails. The agentic UI substrate has to ship them by default."),
    ]
    bullets(s, Inches(0.7), Inches(1.7), W - Inches(1.4), Inches(5.0),
            drivers, size=14, gap=10)
    add_speaker_notes(s,
        "Why this matters now: protocols + tool-calling + MCP + generative UI all matured in a 12-month window.",
        "Buying or building? The protocols change every quarter. A library that abstracts them is the right unit of investment.")
    return s


def slide_three_protocols(prs):
    s = add(prs)
    slide_chrome(s, "1 · THE SHIFT", "Three protocols, one library", "Architects")

    # Three protocol cards in a row
    card_w = Inches(3.9); card_h = Inches(4.2); gap = Inches(0.3)
    start_x = Inches(0.7)
    protocols = [
        ("AG-UI", "CopilotKit",
         RGBColor(0x4F, 0x46, 0xE5),
         RGBColor(0xEE, 0xF2, 0xFF),
         ["Streaming via SSE",
          "Lifecycle + tool-call events",
          "Generative UI via show-component",
          "Mature client implementation",
          "Best for: server-driven agent flows"]),
        ("Hashbrown", "Liquid Frontiers",
         RGBColor(0x05, 0x96, 0x69),
         RGBColor(0xD1, 0xFA, 0xE5),
         ["UI-generation streams",
          "Widget + props natively",
          "Multi-provider (OpenAI, Google)",
          "Lighter than AG-UI",
          "Best for: rapid generative-UI demos"]),
        ("A2UI", "Multi-vendor consortium",
         RGBColor(0xD9, 0x77, 0x06),
         RGBColor(0xFE, 0xF3, 0xC7),
         ["ui-action event class",
          "Agent drives the whole UI",
          "Action-driven (not just tool-driven)",
          "Spec still evolving",
          "Best for: agent-driven nav + state"]),
    ]
    for i, (name, who, fg, bg, body) in enumerate(protocols):
        x = start_x + i * (card_w + gap)
        rounded(s, x, Inches(1.8), card_w, card_h, SURFACE, line=BORDER, radius=0.04)
        rect(s, x, Inches(1.8), card_w, Inches(0.5), fg)
        text_box(s, x + Inches(0.25), Inches(1.92), card_w - Inches(0.5),
                 Inches(0.3), name, size=18, bold=True, color=SURFACE)
        text_box(s, x + Inches(0.25), Inches(2.4), card_w - Inches(0.5),
                 Inches(0.4), f"By {who}", size=11, bold=True, color=fg)
        bullets(s, x + Inches(0.25), Inches(2.85), card_w - Inches(0.5),
                Inches(2.8), body, size=12, gap=6, marker="·",
                marker_color=fg)

    text_box(s, Inches(0.7), Inches(6.2), W - Inches(1.4), Inches(0.8),
             "@maverick/agentic-ui — one AgenticBackend abstraction; chat shell, registries, and widgets stay protocol-agnostic.",
             size=14, bold=True, color=BRAND_DEEP, align=PP_ALIGN.CENTER)
    add_speaker_notes(s,
        "AG-UI: most mature, best for server-driven flows where the agent emits streamed text + tool calls + widgets.",
        "Hashbrown: lighter weight, faster to prototype, multi-provider out of the box.",
        "A2UI: distinguishing feature is the ui-action class — the agent can navigate, fill forms, mutate stores. Spec still moving.")
    return s


# ── AG-UI deep dive ─────────────────────────────────────────────────────────

def slide_agui_background(prs):
    s = add(prs)
    slide_chrome(s, "2 · AG-UI", "Background & origin", "Architects")

    text_box(s, Inches(0.7), Inches(1.7), W - Inches(1.4), Inches(0.5),
             "Open protocol for agent ↔ UI communication, originated by CopilotKit.",
             size=16, color=TEXT_2)

    facts = [
        ("Transport-agnostic.", "SSE in practice; WebSocket and plain HTTP also valid carriers."),
        ("Message-based.", "Discrete event types — RUN_STARTED, TEXT_MESSAGE_*, TOOL_CALL_*, RUN_FINISHED — not a free-form stream."),
        ("Tool-call first-class.", "Distinguishes server-side tools (run on the agent host) from client-side tools (run in the browser)."),
        ("Generative UI by convention.", "Reserved show-component tool whose result tells the host which registered widget to render."),
        ("Reference implementation.", "@ag-ui/client — the runUntilSettled loop + HttpAgent SSE adapter — is the canonical client."),
        ("Adoption.", "CopilotKit, AG-UI Studio, multiple Mastra-based agent backends, several Angular and React reference apps."),
    ]
    bullets(s, Inches(0.7), Inches(2.4), W - Inches(1.4), Inches(4.5),
            facts, size=13, gap=8)
    add_speaker_notes(s,
        "AG-UI is the one most enterprise teams pick today. The reason: it cleanly maps lifecycle, text, and tool calls to discrete events — easy to log, easy to audit.",
        "Origin: CopilotKit. Anglo/architects.io article is the canonical primer.")
    return s


def slide_agui_events(prs):
    s = add(prs)
    slide_chrome(s, "2 · AG-UI", "Event model — the wire format", "Developers")

    rounded(s, Inches(0.7), Inches(1.7), W - Inches(1.4), Inches(5.4),
            SURFACE, line=BORDER, radius=0.02)

    sections = [
        ("Lifecycle",   ["RUN_STARTED", "RUN_FINISHED", "RUN_ERROR"], BRAND),
        ("Text",        ["TEXT_MESSAGE_START", "TEXT_MESSAGE_CONTENT", "TEXT_MESSAGE_END"], INFO),
        ("Tool calls",  ["TOOL_CALL_START", "TOOL_CALL_ARGS", "TOOL_CALL_END", "TOOL_CALL_RESULT"], OK),
        ("Generative",  ["WIDGET_RENDER (synth from show-component)", "UI_ACTION (A2UI extension)"], WARN),
    ]
    y = Inches(1.95)
    for title, evts, color in sections:
        rect(s, Inches(0.9), y, Inches(0.16), Inches(1.0), color)
        text_box(s, Inches(1.2), y, Inches(2.5), Inches(0.4),
                 title, size=14, bold=True, color=color)
        text_box(s, Inches(1.2), y + Inches(0.4), Inches(11), Inches(0.6),
                 "  ·  ".join(evts), size=12, color=TEXT_2,
                 font="Menlo")
        y += Inches(1.15)

    text_box(s, Inches(0.7), Inches(6.7), W - Inches(1.4), Inches(0.4),
             "Mapped 1:1 in @maverick/agentic-ui's AgenticBackend.run() event union — adapters never re-shape the model.",
             size=11, color=MUTED, align=PP_ALIGN.CENTER)
    add_speaker_notes(s,
        "Show this slide to engineers. Each event has a stable schema; AG-UI never embeds free-form payloads at the boundary.",
        "Our adapter @maverick/agentic-ui/ag-ui maps these 1:1 — you can build a debugger by tapping the AgenticBackend stream.")
    return s


def slide_agui_capabilities(prs):
    s = add(prs)
    slide_chrome(s, "2 · AG-UI", "Capabilities & where it fits")

    # Two columns: Strengths and Gaps
    col_w = (W - Inches(1.4) - Inches(0.3)) / 2
    rounded(s, Inches(0.7), Inches(1.7), col_w, Inches(5.0), SURFACE, line=BORDER)
    rect(s, Inches(0.7), Inches(1.7), Inches(0.06), Inches(5.0), OK)
    text_box(s, Inches(0.95), Inches(1.85), col_w - Inches(0.4), Inches(0.4),
             "Strengths", size=15, bold=True, color=OK)
    bullets(s, Inches(0.95), Inches(2.4), col_w - Inches(0.4), Inches(4.4),
            ["Cleanest event model of the three",
             "Server-side tool execution baked in",
             "SSE — proxies and CDNs handle it natively",
             "Mature CopilotKit + Mastra ecosystem",
             "Trace-friendly: each event is a structured record",
             "Best fit for compliance / audit / regulated domains"],
            size=12, gap=6, marker_color=OK)

    rounded(s, Inches(0.7) + col_w + Inches(0.3), Inches(1.7), col_w,
            Inches(5.0), SURFACE, line=BORDER)
    rect(s, Inches(0.7) + col_w + Inches(0.3), Inches(1.7), Inches(0.06),
         Inches(5.0), WARN)
    text_box(s, Inches(0.95) + col_w + Inches(0.3), Inches(1.85),
             col_w - Inches(0.4), Inches(0.4),
             "Trade-offs / gaps", size=15, bold=True, color=WARN)
    bullets(s, Inches(0.95) + col_w + Inches(0.3), Inches(2.4),
            col_w - Inches(0.4), Inches(4.4),
            ["Heavier than Hashbrown — more events to handle",
             "ui-action / agent-driven nav not in core spec",
             "Generative UI is via convention, not a first-class event",
             "No standard form schema — apps roll their own",
             "Per-protocol adapter still required (we provide one)"],
            size=12, gap=6, marker_color=WARN)
    add_speaker_notes(s,
        "When to pick AG-UI: regulated industries, audit-grade trails, server-side tool catalogs.",
        "When NOT: lightweight prototypes, mostly-client-rendered UIs with little server orchestration — Hashbrown is faster to spin up.")
    return s


# ── Hashbrown deep dive ─────────────────────────────────────────────────────

def slide_hashbrown_background(prs):
    s = add(prs)
    slide_chrome(s, "3 · HASHBROWN", "Background & origin", "Architects")

    text_box(s, Inches(0.7), Inches(1.7), W - Inches(1.4), Inches(0.5),
             "Open framework from Liquid Frontiers focused on generative UI and tool-calling.",
             size=16, color=TEXT_2)
    facts = [
        ("UI-generation native.", "Hashbrown's streams emit component + props directly — no show-component shim required."),
        ("Multi-provider out of the box.", "OpenAI, Google, and Anthropic — swap providers via a config flag."),
        ("Smaller surface than AG-UI.", "Fewer event types; lifecycle is implicit in stream open/close."),
        ("React-first heritage.", "Original target was React; Angular adapter is recent. Our wrapper makes it idiomatic for either."),
        ("Schema-driven UI generation.", "Zod-based component schemas anchor the LLM's structured output."),
        ("Best fit.", "Rapid prototypes, demo-heavy generative UI, multi-LLM A/B comparisons."),
    ]
    bullets(s, Inches(0.7), Inches(2.4), W - Inches(1.4), Inches(4.5),
            facts, size=13, gap=8)
    add_speaker_notes(s,
        "Hashbrown's pitch: thinner protocol, faster start, multi-provider built in.",
        "Trade-off: the lifecycle isn't as crisp, so audit-grade logging needs extra work.")
    return s


def slide_hashbrown_features(prs):
    s = add(prs)
    slide_chrome(s, "3 · HASHBROWN", "Features at a glance", "Developers")

    bullet_rows = [
        ("ui-stream",    "Emits widget + props natively — no shim", BRAND),
        ("text-stream",  "Token-level text deltas",                  INFO),
        ("tool-stream",  "Function-call args + results",             OK),
        ("provider-mux", "Switch OpenAI ↔ Google ↔ Anthropic at runtime", WARN),
        ("schema-anchor","Zod schemas anchor structured LLM output",  BAD),
        ("react-bridge", "First-class React; Angular via wrapper",    MUTED),
    ]
    y = Inches(1.85)
    row_h = Inches(0.74)
    for tag, desc, color in bullet_rows:
        rounded(s, Inches(0.7), y, W - Inches(1.4), row_h - Inches(0.1),
                SURFACE, line=BORDER, radius=0.04)
        chip(s, Inches(0.9), y + Inches(0.13), tag, fg=SURFACE, bg=color,
             size=11, w=Inches(1.6))
        text_box(s, Inches(2.7), y + Inches(0.18), W - Inches(3.6),
                 Inches(0.4), desc, size=13, color=TEXT_2)
        y += row_h
    add_speaker_notes(s,
        "Highlight provider-mux — for an org doing LLM evaluation, Hashbrown lets you swap providers without rewiring the chat shell.",
        "Highlight schema-anchor — Zod schemas are the contract; same library + skill in our schematics.")
    return s


def slide_hashbrown_vs_agui(prs):
    s = add(prs)
    slide_chrome(s, "3 · HASHBROWN", "Hashbrown vs AG-UI — when to pick what")

    rows = [
        ("Lifecycle clarity",   "Implicit (stream open/close)", "Explicit (RUN_STARTED/FINISHED)", "AG-UI"),
        ("Generative UI",       "Native (UI-stream)",            "Convention (show-component)",     "Hashbrown"),
        ("Audit-grade logging", "Needs work — synth lifecycle",  "Drop-in — events are records",    "AG-UI"),
        ("Time-to-first-prompt","Hours",                          "Day or two",                      "Hashbrown"),
        ("Multi-provider",      "Built in",                       "Adapter per provider",             "Hashbrown"),
        ("Server-side tools",   "Adapter needed",                 "First-class",                     "AG-UI"),
        ("Spec stability",      "Smaller surface · stable",       "Larger surface · stable",         "Tie"),
    ]
    rounded(s, Inches(0.7), Inches(1.7), W - Inches(1.4), Inches(5.4),
            SURFACE, line=BORDER, radius=0.01)

    # Header
    rect(s, Inches(0.7), Inches(1.7), W - Inches(1.4), Inches(0.5), SURFACE_2)
    text_box(s, Inches(0.9),  Inches(1.78), Inches(3.0), Inches(0.4),
             "Dimension", size=11, bold=True, color=MUTED)
    text_box(s, Inches(4.0),  Inches(1.78), Inches(3.0), Inches(0.4),
             "Hashbrown", size=11, bold=True, color=OK)
    text_box(s, Inches(7.2),  Inches(1.78), Inches(3.0), Inches(0.4),
             "AG-UI", size=11, bold=True, color=BRAND)
    text_box(s, Inches(10.6), Inches(1.78), Inches(2.0), Inches(0.4),
             "Pick", size=11, bold=True, color=TEXT)

    y = Inches(2.2)
    for dim, hb, agui, pick in rows:
        text_box(s, Inches(0.9), y, Inches(3.0), Inches(0.5),
                 dim, size=12, bold=True, color=TEXT)
        text_box(s, Inches(4.0), y, Inches(3.1), Inches(0.5),
                 hb, size=12, color=TEXT_2)
        text_box(s, Inches(7.2), y, Inches(3.3), Inches(0.5),
                 agui, size=12, color=TEXT_2)
        chip(s, Inches(10.6), y + Inches(0.05), pick,
             fg=BRAND_DEEP if pick != "Tie" else MUTED,
             bg=BRAND_TINT if pick != "Tie" else SURFACE_2,
             size=10, w=Inches(1.5))
        y += Inches(0.66)
    add_speaker_notes(s,
        "Decision matrix to share with architects. Not a verdict — a guide.",
        "Most enterprises pick AG-UI for production; Hashbrown for prototyping.")
    return s


# ── A2UI ────────────────────────────────────────────────────────────────────

def slide_a2ui(prs):
    s = add(prs)
    slide_chrome(s, "4 · A2UI", "Agent-driven UI — ui-action as a class", "Architects")

    text_box(s, Inches(0.7), Inches(1.7), W - Inches(1.4), Inches(0.5),
             "A2UI lifts the agent from 'calls tools' to 'drives the whole UI'.",
             size=16, color=TEXT_2)

    text_box(s, Inches(0.7), Inches(2.4), W - Inches(1.4), Inches(0.4),
             "Distinguishing event class", size=14, bold=True, color=WARN)
    code_block(s, Inches(0.7), Inches(2.85), W - Inches(1.4), Inches(1.6), [
        "{ type: 'ui-action',",
        "  actionId: 'a-7f2e',",
        "  op: 'navigate',          // or 'applyTag', 'fillForm', 'openDrawer', …",
        "  payload: { path: '/documents', queryParams: { id: 'DOC-1234' } } }",
    ])

    text_box(s, Inches(0.7), Inches(4.7), W - Inches(1.4), Inches(0.4),
             "Where it shines", size=14, bold=True, color=BRAND)
    bullets(s, Inches(0.7), Inches(5.15), W - Inches(1.4), Inches(2.0), [
        "Agent navigates to a route after a tool call settles",
        "Form-fill flows where the agent populates fields the user verifies",
        "State mutations dispatched into NgRx / Redux without a separate tool",
        "Cross-MFE coordination — one remote raises an action a different remote handles",
    ], size=13, gap=6, marker_color=BRAND)

    add_speaker_notes(s,
        "A2UI is younger; specification still evolving. We reserve ui-action in our event union from M1 so the chat shell never breaks when the spec moves.",
        "Same pattern — clickable widgets in our eDiscovery demo dispatch openCustodian/openDocument/openHold actions through ActionRegistry.")
    return s


# ── @maverick/agentic-ui — the library ──────────────────────────────────────

def slide_lib_problem(prs):
    s = add(prs)
    slide_chrome(s, "5 · THE LIBRARY", "Why a library", "Senior executives")

    text_box(s, Inches(0.7), Inches(1.7), W - Inches(1.4), Inches(0.5),
             "Three protocols, four bundlers, two MFE runtimes, and an MCP transport — that's not a library, that's a forklift.",
             size=15, color=TEXT_2)
    rounded(s, Inches(0.7), Inches(2.5), Inches(6.0), Inches(4.4),
            RGBColor(0xFE, 0xF2, 0xF2), line=BAD, radius=0.04)
    text_box(s, Inches(0.95), Inches(2.65), Inches(5.5), Inches(0.4),
             "What you'd build without it", size=14, bold=True, color=BAD)
    bullets(s, Inches(0.95), Inches(3.15), Inches(5.5), Inches(3.5), [
        "Per-protocol chat shell, three times",
        "Hand-rolled tool registry, ad hoc per app",
        "Bespoke MFE federation glue for each remote",
        "MCP server scaffolding from scratch",
        "Telemetry stitched in late, never quite consistent",
        "Schematics: none — every team retypes the bootstrap",
    ], size=12, gap=6, marker="✗", marker_color=BAD)

    rounded(s, Inches(7.0), Inches(2.5), Inches(5.7), Inches(4.4),
            RGBColor(0xEC, 0xFD, 0xF5), line=OK, radius=0.04)
    text_box(s, Inches(7.25), Inches(2.65), Inches(5.2), Inches(0.4),
             "What @maverick/agentic-ui ships", size=14, bold=True, color=OK)
    bullets(s, Inches(7.25), Inches(3.15), Inches(5.2), Inches(3.5), [
        "One AgenticBackend interface — three adapters",
        "13 typed registries, one base class, signal-backed",
        "Native + Module Federation, both first-class",
        "@maverick/agentic-ui-mcp wraps tools as MCP",
        "OTel-instrumented from M1; default no-op",
        "ng add scaffolds an entire app in one command",
    ], size=12, gap=6, marker="✓", marker_color=OK)
    add_speaker_notes(s,
        "Frame for execs: building from scratch is a 6+ month investment. The library is a 2-week pilot.",
        "Production-grade, MIT-licensed, MIT-friendly. No vendor lock-in to a paid SaaS.")
    return s


def slide_lib_architecture(prs):
    s = add(prs)
    slide_chrome(s, "5 · THE LIBRARY", "Architecture at a glance", "Architects")

    # ── Browser host container ─────────────────────────────────────────────
    host_x, host_y = Inches(0.7), Inches(1.7)
    host_w, host_h = Inches(8.5), Inches(5.0)
    rounded(s, host_x, host_y, host_w, host_h, SURFACE_2, line=BORDER, radius=0.02)
    text_box(s, host_x + Inches(0.2), host_y + Inches(0.1), Inches(4),
             Inches(0.3), "BROWSER HOST", size=10, bold=True, color=MUTED)

    # UI layer (top of host)
    ui_y = host_y + Inches(0.5)
    ui_w = Inches(2.4); ui_h = Inches(0.6); gap = Inches(0.15)
    chip_node(s, host_x + Inches(0.3), ui_y, ui_w, ui_h,
              "<mvk-chat-shell>", fill=BRAND, fg=SURFACE)
    chip_node(s, host_x + Inches(0.3) + ui_w + gap, ui_y, ui_w, ui_h,
              "Sidebar nav", fill=BRAND, fg=SURFACE)
    chip_node(s, host_x + Inches(0.3) + 2 * (ui_w + gap), ui_y, ui_w, ui_h,
              "Routed pages", fill=BRAND, fg=SURFACE)

    # Registry hub
    reg_y = ui_y + ui_h + Inches(0.45)
    reg_w = host_w - Inches(0.6); reg_h = Inches(1.2)
    rounded(s, host_x + Inches(0.3), reg_y, reg_w, reg_h,
            BRAND_TINT, line=BRAND_DEEP, radius=0.03)
    text_box(s, host_x + Inches(0.4), reg_y + Inches(0.1),
             reg_w - Inches(0.2), Inches(0.3),
             "REGISTRY LAYER · 13 typed registries · signal-backed", size=10,
             bold=True, color=BRAND_DEEP)
    # Inner registry chips
    chips = [("Tool", BRAND), ("Component", BRAND), ("Capability", INFO),
             ("Backend", INFO), ("Action", OK), ("Form", OK),
             ("Intent", OK), ("DataSource", WARN), ("Validation", WARN),
             ("Persistence", BAD), ("Layout", BAD), ("Schema", BAD)]
    chip_w = (reg_w - Inches(0.4)) / 6
    for i, (lbl, c) in enumerate(chips):
        col = i % 6; row = i // 6
        cx = host_x + Inches(0.4) + col * chip_w
        cy = reg_y + Inches(0.45) + row * Inches(0.32)
        chip_node(s, cx, cy, chip_w - Inches(0.05), Inches(0.28),
                  lbl, fill=SURFACE, fg=c, border=BORDER, size=9)

    # Backend adapters (bottom of host)
    ad_y = reg_y + reg_h + Inches(0.45)
    ad_w = (host_w - Inches(0.9)) / 3
    backends = [("AgUiBackend", BRAND), ("HashbrownBackend", OK), ("A2uiBackend", WARN)]
    for i, (name, c) in enumerate(backends):
        ax = host_x + Inches(0.3) + i * (ad_w + Inches(0.15))
        chip_node(s, ax, ad_y, ad_w, Inches(0.5), name,
                  fill=SURFACE, fg=c, border=c, size=11)

    # Footer caption inside host
    text_box(s, host_x + Inches(0.3), ad_y + Inches(0.7), host_w - Inches(0.6),
             Inches(0.4),
             "All pushed events also flow into AgenticTelemetrySink / AgenticLogger.",
             size=9, color=MUTED, align=PP_ALIGN.CENTER)

    # ── Right column: external systems ─────────────────────────────────────
    rx = host_x + host_w + Inches(0.3)
    rw = W - rx - Inches(0.7)
    # Federation runtime
    fed_y = host_y
    rounded(s, rx, fed_y, rw, Inches(1.6), SURFACE, line=BORDER, radius=0.04)
    rect(s, rx, fed_y, Inches(0.06), Inches(1.6), BRAND_DEEP)
    text_box(s, rx + Inches(0.2), fed_y + Inches(0.15), rw - Inches(0.3),
             Inches(0.3), "Federation runtime", size=11, bold=True, color=BRAND_DEEP)
    text_box(s, rx + Inches(0.2), fed_y + Inches(0.5), rw - Inches(0.3),
             Inches(1.0),
             "Native Federation\n· esbuild-native\n· @angular-architects\n\n"
             "Module Federation\n· webpack peer", size=9, color=TEXT_2)

    # Agent server
    ag_y = fed_y + Inches(1.75)
    rounded(s, rx, ag_y, rw, Inches(1.4), SURFACE, line=BORDER, radius=0.04)
    rect(s, rx, ag_y, Inches(0.06), Inches(1.4), OK)
    text_box(s, rx + Inches(0.2), ag_y + Inches(0.15), rw - Inches(0.3),
             Inches(0.3), "Agent server", size=11, bold=True, color=OK)
    text_box(s, rx + Inches(0.2), ag_y + Inches(0.5), rw - Inches(0.3),
             Inches(0.9),
             "Mastra · Gemini · OpenAI\nAG-UI SSE route\nThreadStateStore",
             size=9, color=TEXT_2)

    # MCP server
    mcp_y = ag_y + Inches(1.55)
    rounded(s, rx, mcp_y, rw, Inches(1.1), SURFACE, line=BORDER, radius=0.04)
    rect(s, rx, mcp_y, Inches(0.06), Inches(1.1), WARN)
    text_box(s, rx + Inches(0.2), mcp_y + Inches(0.1), rw - Inches(0.3),
             Inches(0.3), "MCP server", size=11, bold=True, color=WARN)
    text_box(s, rx + Inches(0.2), mcp_y + Inches(0.42), rw - Inches(0.3),
             Inches(0.7),
             "Same ToolDef\nClaude Desktop / Cursor\nText/html;profile=mcp-app",
             size=9, color=TEXT_2)

    # Arrows: backend → external
    arrow(s, host_x + Inches(0.3) + ad_w + Inches(0.07),
          ad_y + Inches(0.25),
          rx, ag_y + Inches(0.7), color=OK, weight=1.5)
    arrow(s, host_x + Inches(0.3) + 2 * (ad_w + Inches(0.15)) + ad_w / 2,
          ad_y + Inches(0.5),
          rx, mcp_y + Inches(0.55), color=WARN, weight=1.5, dashed=True)
    # Arrow: federation → registry
    arrow(s, rx, fed_y + Inches(0.8),
          host_x + host_w - Inches(0.05), reg_y + Inches(0.6),
          color=BRAND_DEEP, weight=1.5)

    # Bottom caption
    text_box(s, Inches(0.7), Inches(7.05), W - Inches(1.4), Inches(0.3),
             "One library · three protocols · two federation paths · MCP for analyst desktops",
             size=11, bold=True, color=BRAND_DEEP, align=PP_ALIGN.CENTER)

    add_speaker_notes(s,
        "Topology view. The host browser is the centre — registries are the only shared state.",
        "Federation runtime feeds remotes' capabilities INTO the registry layer; backend adapters ship runs OUT to the agent server / MCP.",
        "Telemetry crosscuts every layer; the OTel exporter ships in /otel and is opt-in.")
    return s


def slide_capability_handoff(prs):
    """NEW: sequence diagram of MFE capability handoff."""
    s = add(prs)
    slide_chrome(s, "5 · THE LIBRARY", "Capability handoff — sequence diagram", "Architects")

    # Five lifelines
    actors = [
        ("Host shell",  Inches(1.7),  BRAND),
        ("MfeRegistry", Inches(4.0),  INFO),
        ("Federation",  Inches(6.4),  BRAND_DEEP),
        ("Remote",      Inches(8.8),  OK),
        ("Registries",  Inches(11.4), WARN),
    ]
    top = Inches(1.85); bot = Inches(6.6)
    for label, x, color in actors:
        actor_lifeline(s, x, top, bot, label, color=color)

    # Messages — y rows for each step
    rows = [
        (Inches(2.65), 0, 1, "discover(env)", BRAND),
        (Inches(3.05), 1, 0, "RemoteSpec[]", INFO),
        (Inches(3.55), 0, 2, "loadRemoteCapabilities({remote, loader})", BRAND_DEEP),
        (Inches(3.95), 2, 3, "import './Capability'", BRAND_DEEP),
        (Inches(4.35), 3, 4, "module.apply(injector)", OK),
        (Inches(4.85), 4, 4, "registerAll(tools, components)", WARN),
        (Inches(5.35), 0, 4, "next chat turn → ToolRegistry.signal()", BRAND),
        (Inches(5.85), 0, 0, "user prompt fires; agent sees new tools", BRAND),
    ]
    for y, src, dst, label, color in rows:
        x1 = actors[src][1]; x2 = actors[dst][1]
        if src == dst:
            # self-call: short loop
            arrow(s, x1, y, x1 + Inches(0.6), y, color=color, weight=1.2)
            arrow(s, x1 + Inches(0.6), y, x1 + Inches(0.6), y + Inches(0.2),
                  color=color, weight=1.2)
            arrow(s, x1 + Inches(0.6), y + Inches(0.2), x1 + Inches(0.05),
                  y + Inches(0.2), color=color, weight=1.2)
            text_box(s, x1 + Inches(0.7), y - Inches(0.05),
                     Inches(3), Inches(0.3), label, size=9, color=color, bold=True)
        else:
            arrow(s, x1, y, x2, y, color=color, weight=1.2)
            mid_x = (x1 + x2) / 2
            text_box(s, mid_x - Inches(2), y - Inches(0.32),
                     Inches(4), Inches(0.3), label, size=9, color=color,
                     bold=True, align=PP_ALIGN.CENTER)

    # Caption
    text_box(s, Inches(0.7), Inches(6.85), W - Inches(1.4), Inches(0.4),
             "Each remote is another row in this sequence — registries are the single seam every remote writes through.",
             size=11, color=MUTED, align=PP_ALIGN.CENTER)
    add_speaker_notes(s,
        "Sequence diagram for federation handoff. Five actors, eight messages.",
        "Note the last message — ToolRegistry's signal notifies the chat shell on the next turn, mid-session, with no reload.")
    return s


def slide_production_chain(prs):
    """NEW: Phase 3 production chain — sequence flow."""
    s = add(prs)
    slide_chrome(s, "6 · EXAMPLES", "eDiscovery — production chain (Phase 3)",
                 "Developers")

    text_box(s, Inches(0.7), Inches(1.65), W - Inches(1.4), Inches(0.4),
             "One prompt drives a 4-step chained tool call · each step renders its own widget · ActionRegistry navigates on click.",
             size=12, color=TEXT_2)

    # Four numbered steps in a row
    step_w = (W - Inches(1.4) - Inches(0.6)) / 4
    step_h = Inches(2.3)
    step_y = Inches(2.2)
    steps = [
        (1, "createProductionSet", "scope filter · format · Bates pattern\nValidation: pattern conformance",
         "→ productionSummary widget\n   status: draft", BRAND),
        (2, "redactDocument*", "page · bbox · reason\n(* optional, repeatable)",
         "→ redactionEditor widget\n   canvas overlay", BRAND_DEEP),
        (3, "assignBatesNumbers", "stamp sequential ids per pattern\nstart = 1 (default)",
         "→ batesPreview + summary\n   status: review", WARN),
        (4, "exportProductionSet", "deliver: true (irreversible)\nrequires audit reason",
         "→ productionSummary\n   status: delivered", OK),
    ]
    for i, (n, name, body, returns, color) in enumerate(steps):
        x = Inches(0.7) + i * (step_w + Inches(0.2))
        rounded(s, x, step_y, step_w, step_h, SURFACE, line=BORDER, radius=0.04)
        rect(s, x, step_y, step_w, Inches(0.05), color)
        # Numbered badge
        rounded(s, x + Inches(0.2), step_y + Inches(0.18), Inches(0.4),
                Inches(0.4), color, radius=0.5)
        text_box(s, x + Inches(0.2), step_y + Inches(0.18), Inches(0.4),
                 Inches(0.4), str(n), size=14, bold=True, color=SURFACE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x + Inches(0.7), step_y + Inches(0.2),
                 step_w - Inches(0.85), Inches(0.32), name,
                 size=11, bold=True, color=color, font="Menlo")
        text_box(s, x + Inches(0.2), step_y + Inches(0.7),
                 step_w - Inches(0.4), Inches(0.7),
                 body, size=10, color=TEXT_2)
        # Returns block
        rounded(s, x + Inches(0.2), step_y + Inches(1.55),
                step_w - Inches(0.4), Inches(0.65), SURFACE_2,
                line=BORDER, radius=0.03)
        text_box(s, x + Inches(0.3), step_y + Inches(1.6),
                 step_w - Inches(0.5), Inches(0.6),
                 returns, size=9, color=color, bold=True)
        # Connector arrow to next step
        if i < 3:
            arrow(s, x + step_w, step_y + step_h / 2,
                  x + step_w + Inches(0.18), step_y + step_h / 2,
                  color=color, weight=2)

    # Side band: artefacts touched
    band_y = step_y + step_h + Inches(0.3)
    band_h = Inches(1.4)
    rounded(s, Inches(0.7), band_y, W - Inches(1.4), band_h, SURFACE_2,
            line=BORDER, radius=0.02)
    text_box(s, Inches(0.95), band_y + Inches(0.12), Inches(4), Inches(0.3),
             "Library seams exercised", size=11, bold=True, color=BRAND_DEEP)

    seams = [
        ("ToolRegistry", "4 tools registered via federation", BRAND),
        ("ComponentRegistry", "3 widgets contributed by remote", BRAND),
        ("FormRegistry", "productionConfigForm — schema-driven", OK),
        ("ValidationRegistry", "Bates pattern conformance check", WARN),
        ("ActionRegistry", "openProduction — click-to-navigate", BRAND_DEEP),
        ("Audit log", "every step appends an event", BAD),
    ]
    chip_x = Inches(0.95)
    chip_y = band_y + Inches(0.55)
    for i, (name, desc, c) in enumerate(seams):
        col = i % 3; row = i // 3
        cx = Inches(0.95) + col * Inches(4.1)
        cy = band_y + Inches(0.5) + row * Inches(0.35)
        chip_node(s, cx, cy, Inches(1.5), Inches(0.28), name,
                  fill=SURFACE, fg=c, border=c, size=9)
        text_box(s, cx + Inches(1.55), cy + Inches(0.04),
                 Inches(2.5), Inches(0.3), desc, size=9, color=TEXT_2)

    add_speaker_notes(s,
        "This is THE slide for the eDiscovery case study — shows the chained-tool-call pattern with widget feedback at every step.",
        "Six seams in one workflow. That's why we needed registries — each seam is independently extensible.")
    return s


def slide_full_picture(prs):
    """NEW v4: complete-picture architecture diagram. Shows every layer
    end-to-end with directional communication arrows so the audience
    can trace one user prompt from the browser to the LLM and back."""
    s = add(prs)
    slide_chrome(s, "05 · ARCHITECTURE",
                 "End-to-end picture — one prompt, all layers",
                 "Architects")

    text_box(s, Inches(0.7), Inches(1.85), W - Inches(1.4), Inches(0.4),
             "User types a prompt → arrows show every hop until widgets land back in the chat.",
             size=11, color=MUTED)

    # Three vertical bands: Browser  |  Server  |  External
    bx = Inches(0.7); by = Inches(2.4)
    band_w = (W - Inches(1.4)) / 3 - Inches(0.1)
    band_h = Inches(4.4)

    # Band 1 — Browser host
    rounded(s, bx, by, band_w, band_h, SURFACE, line=BORDER, radius=0.02)
    text_box(s, bx + Inches(0.15), by + Inches(0.1), band_w - Inches(0.3),
             Inches(0.3), "BROWSER · ANGULAR HOST", size=9, bold=True, color=MUTED)

    # User node
    chip_node(s, bx + Inches(0.4), by + Inches(0.55), band_w - Inches(0.8),
              Inches(0.45), "User prompt",
              fill=ACCENT_TINT, fg=ACCENT, border=ACCENT, size=11)
    # Chat shell
    chip_node(s, bx + Inches(0.4), by + Inches(1.15), band_w - Inches(0.8),
              Inches(0.45), "<mvk-chat-shell>",
              fill=BRAND, fg=SURFACE, size=11)
    # injectAgenticChat / runUntilSettled
    chip_node(s, bx + Inches(0.4), by + Inches(1.75), band_w - Inches(0.8),
              Inches(0.45), "injectAgenticChat()",
              fill=BRAND_DEEP, fg=SURFACE, size=10)
    # Tool filter
    chip_node(s, bx + Inches(0.4), by + Inches(2.35), band_w - Inches(0.8),
              Inches(0.45), "keywordToolFilter (top-12)",
              fill=ACCENT, fg=SURFACE, size=10)
    # Backend adapter
    chip_node(s, bx + Inches(0.4), by + Inches(2.95), band_w - Inches(0.8),
              Inches(0.45), "AgUiBackend",
              fill=BRAND, fg=SURFACE, size=10)
    # ToolRegistry / ComponentRegistry
    chip_node(s, bx + Inches(0.4), by + Inches(3.55), band_w - Inches(0.8),
              Inches(0.45), "ToolRegistry · ComponentRegistry",
              fill=INFO, fg=SURFACE, size=9)

    # Band 2 — Server
    bx2 = bx + band_w + Inches(0.15)
    rounded(s, bx2, by, band_w, band_h, SURFACE, line=BORDER, radius=0.02)
    text_box(s, bx2 + Inches(0.15), by + Inches(0.1), band_w - Inches(0.3),
             Inches(0.3), "AGENT SERVER · NODE", size=9, bold=True, color=MUTED)

    chip_node(s, bx2 + Inches(0.4), by + Inches(0.85), band_w - Inches(0.8),
              Inches(0.45), "/agents/coordinator/run",
              fill=BRAND, fg=SURFACE, size=10)
    chip_node(s, bx2 + Inches(0.4), by + Inches(1.45), band_w - Inches(0.8),
              Inches(0.45), "OrchestratorAgent (router)",
              fill=BRAND_DEEP, fg=SURFACE, size=10)
    chip_node(s, bx2 + Inches(0.4), by + Inches(2.05), band_w - Inches(0.8),
              Inches(0.45), "Specialist (review/production/…)",
              fill=BRAND_DEEP, fg=SURFACE, size=9)
    chip_node(s, bx2 + Inches(0.4), by + Inches(2.65), band_w - Inches(0.8),
              Inches(0.45), "GeminiAgent · system prompt",
              fill=INFO, fg=SURFACE, size=10)
    chip_node(s, bx2 + Inches(0.4), by + Inches(3.25), band_w - Inches(0.8),
              Inches(0.45), "ThreadStateStore (sticky)",
              fill=ACCENT, fg=SURFACE, size=10)

    # Band 3 — LLM provider
    bx3 = bx2 + band_w + Inches(0.15)
    rounded(s, bx3, by, band_w, band_h, SURFACE, line=BORDER, radius=0.02)
    text_box(s, bx3 + Inches(0.15), by + Inches(0.1), band_w - Inches(0.3),
             Inches(0.3), "LLM PROVIDER · GEMINI", size=9, bold=True, color=MUTED)

    chip_node(s, bx3 + Inches(0.4), by + Inches(1.0), band_w - Inches(0.8),
              Inches(0.45), "gemini-2.5-flash",
              fill=BRAND, fg=SURFACE, size=11)
    chip_node(s, bx3 + Inches(0.4), by + Inches(1.7), band_w - Inches(0.8),
              Inches(0.45), "Tool-call schema",
              fill=BRAND_DEEP, fg=SURFACE, size=10)
    chip_node(s, bx3 + Inches(0.4), by + Inches(2.4), band_w - Inches(0.8),
              Inches(0.45), "Streamed completion",
              fill=BRAND_DEEP, fg=SURFACE, size=10)
    chip_node(s, bx3 + Inches(0.4), by + Inches(3.1), band_w - Inches(0.8),
              Inches(0.45), "Generative-UI components",
              fill=ACCENT, fg=SURFACE, size=10)

    # Arrows — request path (down-and-across)
    # User → ChatShell
    arrow(s, bx + band_w / 2, by + Inches(1.0),
          bx + band_w / 2, by + Inches(1.15),
          color=ACCENT, weight=2)
    # ChatShell → injectAgenticChat
    arrow(s, bx + band_w / 2, by + Inches(1.6),
          bx + band_w / 2, by + Inches(1.75),
          color=BRAND, weight=1.5)
    # → ToolFilter
    arrow(s, bx + band_w / 2, by + Inches(2.2),
          bx + band_w / 2, by + Inches(2.35),
          color=BRAND, weight=1.5)
    # ToolFilter ↔ ToolRegistry (read)
    arrow(s, bx + band_w / 2 + Inches(0.4), by + Inches(2.575),
          bx + band_w / 2 + Inches(0.4), by + Inches(3.55),
          color=INFO, weight=1, dashed=True)
    # → AgUiBackend
    arrow(s, bx + band_w / 2, by + Inches(2.8),
          bx + band_w / 2, by + Inches(2.95),
          color=BRAND, weight=1.5)
    # AgUiBackend → server route
    arrow(s, bx + band_w - Inches(0.05), by + Inches(3.175),
          bx2 + Inches(0.4), by + Inches(0.95 + 0.225),
          color=BRAND, weight=2)
    # Server → orchestrator → specialist → Gemini chain
    arrow(s, bx2 + band_w / 2, by + Inches(1.3),
          bx2 + band_w / 2, by + Inches(1.45),
          color=BRAND_DEEP, weight=1.5)
    arrow(s, bx2 + band_w / 2, by + Inches(1.9),
          bx2 + band_w / 2, by + Inches(2.05),
          color=BRAND_DEEP, weight=1.5)
    arrow(s, bx2 + band_w / 2, by + Inches(2.5),
          bx2 + band_w / 2, by + Inches(2.65),
          color=INFO, weight=1.5)
    # Specialist → GeminiAPI
    arrow(s, bx2 + band_w - Inches(0.05), by + Inches(2.875),
          bx3 + Inches(0.4), by + Inches(1.075),
          color=INFO, weight=2)

    # Return path (right-to-left, dashed) — LLM → server → browser
    arrow(s, bx3 + Inches(0.4), by + Inches(3.325),
          bx2 + band_w - Inches(0.05), by + Inches(2.875),
          color=ACCENT, weight=1.5, dashed=True)
    arrow(s, bx2 + Inches(0.4), by + Inches(2.275),
          bx + band_w - Inches(0.05), by + Inches(3.175),
          color=ACCENT, weight=1.5, dashed=True)

    # Legend
    ly = by + band_h + Inches(0.15)
    rounded(s, Inches(0.7), ly, W - Inches(1.4), Inches(0.45),
            SURFACE_2, line=BORDER, radius=0.02)
    text_box(s, Inches(0.95), ly + Inches(0.13), Inches(2), Inches(0.25),
             "LEGEND", size=9, bold=True, color=MUTED)
    # Solid arrows
    rect(s, Inches(2.2), ly + Inches(0.22), Inches(0.5), Inches(0.025), BRAND)
    text_box(s, Inches(2.8), ly + Inches(0.13), Inches(2.5), Inches(0.3),
             "Request path", size=10, color=TEXT_2)
    # Dashed arrows
    rect(s, Inches(5.0), ly + Inches(0.22), Inches(0.15), Inches(0.025), ACCENT)
    rect(s, Inches(5.2), ly + Inches(0.22), Inches(0.15), Inches(0.025), ACCENT)
    rect(s, Inches(5.4), ly + Inches(0.22), Inches(0.15), Inches(0.025), ACCENT)
    text_box(s, Inches(5.65), ly + Inches(0.13), Inches(2.5), Inches(0.3),
             "Response · widgets · text deltas", size=10, color=TEXT_2)
    # Tool registry read
    rect(s, Inches(9.0), ly + Inches(0.22), Inches(0.15), Inches(0.025), INFO)
    rect(s, Inches(9.2), ly + Inches(0.22), Inches(0.15), Inches(0.025), INFO)
    text_box(s, Inches(9.45), ly + Inches(0.13), Inches(3.0), Inches(0.3),
             "Registry read · DI cross-cut", size=10, color=TEXT_2)

    add_speaker_notes(s,
        "Single-slide architecture summary. Trace one prompt: User → ChatShell → ToolFilter (reads ToolRegistry) → AgUiBackend → server route → Orchestrator → Specialist → Gemini → events stream back.",
        "Use this slide as the framing for the rest of the architecture deck.")
    return s


def slide_protocol_picker(prs):
    """NEW v4: AG-UI vs Hashbrown — when and where, use cases.
    Anchored on the eDiscovery demo's actual choice (AG-UI)."""
    s = add(prs)
    slide_chrome(s, "02–03 · PROTOCOLS",
                 "AG-UI vs Hashbrown — when and where",
                 "Architects")

    # Top callout — what the eDiscovery demo uses today
    rounded(s, Inches(0.7), Inches(1.8), W - Inches(1.4), Inches(0.85),
            BRAND_TINT, line=BORDER, radius=0.02)
    rect(s, Inches(0.7), Inches(1.8), Inches(0.05), Inches(0.85), BRAND)
    text_box(s, Inches(0.95), Inches(1.92), Inches(2.0), Inches(0.3),
             "THIS DEMO USES", size=9, bold=True, color=BRAND)
    text_box(s, Inches(0.95), Inches(2.18), Inches(11), Inches(0.5),
             "AG-UI for the eDiscovery reference — picked for streaming SSE, "
             "explicit lifecycle events, and audit-grade event records.",
             size=14, color=TEXT)

    # Two-column comparison
    col_w = (W - Inches(1.4) - Inches(0.4)) / 2
    cy = Inches(2.95)

    # AG-UI column
    rounded(s, Inches(0.7), cy, col_w, Inches(3.6),
            SURFACE, line=BORDER, radius=0.02)
    rect(s, Inches(0.7), cy, Inches(0.05), Inches(3.6), BRAND)
    text_box(s, Inches(0.95), cy + Inches(0.15), Inches(3), Inches(0.3),
             "AG-UI", size=18, bold=True, color=BRAND)
    text_box(s, Inches(0.95), cy + Inches(0.55), col_w - Inches(0.5), Inches(0.3),
             "by CopilotKit · streaming · lifecycle-typed", size=10,
             color=MUTED)
    text_box(s, Inches(0.95), cy + Inches(0.95), col_w - Inches(0.5), Inches(0.3),
             "PICK WHEN…", size=9, bold=True, color=BRAND)
    bullets(s, Inches(0.95), cy + Inches(1.25), col_w - Inches(0.5), Inches(2.3), [
        "Audit and compliance bars (legal-tech, healthcare, finance)",
        "Server-side tool catalog managed by a separate team",
        "Mature multi-agent orchestrator (sticky routing)",
        "SSE traverses your existing CDN / proxy stack",
        "Need to trace runs end-to-end through OTel",
    ], size=11, gap=6, marker_color=BRAND)

    # Hashbrown column
    bx = Inches(0.7) + col_w + Inches(0.4)
    rounded(s, bx, cy, col_w, Inches(3.6),
            SURFACE, line=BORDER, radius=0.02)
    rect(s, bx, cy, Inches(0.05), Inches(3.6), ACCENT)
    text_box(s, bx + Inches(0.25), cy + Inches(0.15), Inches(3), Inches(0.3),
             "HASHBROWN", size=18, bold=True, color=ACCENT)
    text_box(s, bx + Inches(0.25), cy + Inches(0.55), col_w - Inches(0.5),
             Inches(0.3),
             "by Liquid Frontiers · ui-stream native · multi-provider",
             size=10, color=MUTED)
    text_box(s, bx + Inches(0.25), cy + Inches(0.95), col_w - Inches(0.5),
             Inches(0.3),
             "PICK WHEN…", size=9, bold=True, color=ACCENT)
    bullets(s, bx + Inches(0.25), cy + Inches(1.25), col_w - Inches(0.5),
            Inches(2.3), [
        "Rapid prototyping — generative UI in a day",
        "Multi-provider eval (OpenAI / Anthropic / Google switch)",
        "UI-generation is the dominant agent output (not tool calls)",
        "Lighter event surface — fewer adapter checks",
        "Audit grade can be deferred to a later sprint",
    ], size=11, gap=6, marker_color=ACCENT)

    # Bottom note — both via the same AgenticBackend
    text_box(s, Inches(0.7), H - Inches(0.85), W - Inches(1.4), Inches(0.3),
             "Both implement the same `AgenticBackend` interface — swapping protocols leaves the chat shell, registries, and widgets untouched.",
             size=11, bold=True, color=BRAND, align=PP_ALIGN.CENTER)

    add_speaker_notes(s,
        "Direct answer to the most common architecture question.",
        "We picked AG-UI for the eDiscovery demo — audit grade matters here. A retail chatbot demo would pick Hashbrown.",
        "The library lets a team start with one and migrate later — no chat-shell rewrite.")
    return s


def slide_run_sequence_step1(prs):
    """NEW v4: step-reveal sequence (frame 1 of 4)."""
    return _build_run_sequence_slide(prs, frame=1)


def slide_run_sequence_step2(prs):
    return _build_run_sequence_slide(prs, frame=2)


def slide_run_sequence_step3(prs):
    return _build_run_sequence_slide(prs, frame=3)


def slide_run_sequence_step4(prs):
    return _build_run_sequence_slide(prs, frame=4)


def _build_run_sequence_slide(prs, *, frame: int):
    """Step-reveal sequence — one slide per frame so PowerPoint's
    auto-advance / click-through gives the appearance of animation
    without relying on python-pptx's flaky animation XML support.

    Frame 1: user prompt fires
    Frame 2: orchestrator routes to specialist
    Frame 3: specialist calls tool, tool result lands
    Frame 4: widget renders + audit event chains
    """
    s = add(prs)
    slide_chrome(s, "06 · CASE STUDY",
                 f"Run sequence — step {frame} of 4",
                 "Developers")

    text_box(s, Inches(0.7), Inches(1.85), W - Inches(1.4), Inches(0.4),
             "Click through frames 1 → 4 to see one chat turn play out across "
             "five actors. Each frame highlights the active arrow.",
             size=11, color=MUTED)

    # Five actors with vertical lifelines
    actor_y_top = Inches(2.55); actor_y_bot = Inches(6.45)
    actors = [
        ("User",         Inches(1.7),  ACCENT,  ['user']),
        ("ChatShell",    Inches(4.0),  BRAND,   ['chat-shell']),
        ("Coordinator",  Inches(6.4),  BRAND,   ['coordinator']),
        ("Specialist",   Inches(8.7),  INFO,    ['specialist']),
        ("ToolRegistry", Inches(11.5), BRAND_DEEP, ['registry']),
    ]
    for label, x, color, _ in actors:
        actor_lifeline(s, x, actor_y_top, actor_y_bot, label, color=color)

    # Messages — each tagged with the frame at which it appears.
    rows = [
        # frame 1 — request leaves user
        (Inches(2.95), 0, 1, "prompt", ACCENT, 1),
        (Inches(3.4),  1, 4, "filter ToolRegistry → top 12", BRAND_DEEP, 1),
        (Inches(3.85), 4, 1, "12 ToolDefs", BRAND_DEEP, 1),
        # frame 2 — chat shell sends to coordinator + routing
        (Inches(4.3),  1, 2, "POST /agents/coordinator/run (SSE)", BRAND, 2),
        (Inches(4.75), 2, 3, "classify → 'production' specialist", BRAND, 2),
        # frame 3 — specialist calls tool
        (Inches(5.2),  3, 1, "tool_call: generateChainOfCustodyReport", INFO, 3),
        (Inches(5.65), 1, 1, "execute handler → audit hash chained", INFO, 3),
        # frame 4 — result + widget
        (Inches(6.1),  1, 3, "tool_result", ACCENT, 4),
        (Inches(6.55), 3, 1, "widget_render: chainOfCustodyReport", ACCENT, 4),
        (Inches(7.0),  1, 0, "rendered card visible", ACCENT, 4),
    ]
    for y, src, dst, label, color, msg_frame in rows:
        active = msg_frame <= frame
        msg_color = color if active else FAINT
        weight = 1.5 if active else 1.0
        x1 = actors[src][1]; x2 = actors[dst][1]
        if src == dst:
            arrow(s, x1, y, x1 + Inches(0.6), y, color=msg_color, weight=weight)
            arrow(s, x1 + Inches(0.6), y, x1 + Inches(0.6), y + Inches(0.18),
                  color=msg_color, weight=weight)
            arrow(s, x1 + Inches(0.6), y + Inches(0.18), x1 + Inches(0.05),
                  y + Inches(0.18), color=msg_color, weight=weight)
            text_box(s, x1 + Inches(0.7), y - Inches(0.05),
                     Inches(4), Inches(0.3), label, size=9,
                     color=msg_color if active else FAINT, bold=active)
        else:
            arrow(s, x1, y, x2, y, color=msg_color, weight=weight,
                  dashed=not active)
            mid_x = (x1 + x2) / 2
            text_box(s, mid_x - Inches(2.2), y - Inches(0.32),
                     Inches(4.4), Inches(0.3), label, size=9,
                     color=msg_color if active else FAINT,
                     bold=active, align=PP_ALIGN.CENTER)

    # Frame caption (highlights what's new this frame)
    captions = {
        1: "Frame 1 — user prompt fires. Tool filter scores 17 tools, sends top 12 to Gemini.",
        2: "Frame 2 — orchestrator classifies and routes the run to the production specialist.",
        3: "Frame 3 — specialist returns a tool call. Handler runs; audit event lands in the chain.",
        4: "Frame 4 — widget renders below the assistant message. Round-trip complete.",
    }
    rounded(s, Inches(0.7), H - Inches(0.95), W - Inches(1.4),
            Inches(0.45), ACCENT_TINT, line=ACCENT, radius=0.04)
    text_box(s, Inches(0.95), H - Inches(0.85), W - Inches(1.7),
             Inches(0.3), captions[frame], size=12, bold=True, color=BRAND_DEEP)

    add_speaker_notes(s,
        f"Step {frame} of 4. Click forward to advance. The four slides build up the same diagram so the audience sees the message flow accumulate.",
        "Use F5 to start full-screen → space-bar to advance.")
    return s


def slide_one_handler_three_surfaces(prs):
    """NEW: 'one handler, multiple surfaces' visualisation."""
    s = add(prs)
    slide_chrome(s, "5 · THE LIBRARY", "One handler, three surfaces", "Architects")

    text_box(s, Inches(0.7), Inches(1.65), W - Inches(1.4), Inches(0.5),
             "The same ToolDef literal powers the chat shell, the MCP server (Claude Desktop / Cursor), and the standalone-mode UI inside each remote — without any code duplication.",
             size=13, color=TEXT_2)

    # Central handler
    cx, cy = W / 2, Inches(4.4)
    handler_w, handler_h = Inches(3.6), Inches(1.3)
    rounded(s, cx - handler_w / 2, cy - handler_h / 2, handler_w, handler_h,
            BRAND, radius=0.04)
    text_box(s, cx - handler_w / 2, cy - handler_h / 2, handler_w, Inches(0.4),
             "ToolDef", size=10, bold=True, color=BRAND_TINT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, cx - handler_w / 2, cy - handler_h / 2 + Inches(0.4),
             handler_w, Inches(0.5),
             "agenticTool({ name, schema, handler })",
             size=12, bold=True, color=SURFACE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font="Menlo")
    text_box(s, cx - handler_w / 2, cy - handler_h / 2 + Inches(0.85),
             handler_w, Inches(0.4),
             "single source of truth · framework-agnostic",
             size=10, color=BRAND_TINT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Three surfaces
    surfaces = [
        ("<mvk-chat-shell>", "Browser · Angular host\nuser ↔ LLM in-app",
         BRAND_DEEP, Inches(2.2), Inches(2.6)),
        ("MCP server", "Claude Desktop · Cursor · Zed\nanalyst workstations",
         WARN, cx - Inches(1.1), Inches(2.6)),
        ("Standalone UI", "Visit :4302 / :4303 directly\ntechnical sanity check",
         OK, W - Inches(4.4), Inches(2.6)),
    ]
    for label, desc, color, x, y in surfaces:
        rounded(s, x, y, Inches(2.2), Inches(1.4), SURFACE,
                line=color, radius=0.04)
        rect(s, x, y, Inches(2.2), Inches(0.05), color)
        text_box(s, x + Inches(0.1), y + Inches(0.15), Inches(2.0),
                 Inches(0.4), label, size=12, bold=True, color=color,
                 align=PP_ALIGN.CENTER)
        text_box(s, x + Inches(0.1), y + Inches(0.55), Inches(2.0),
                 Inches(0.8), desc, size=10, color=TEXT_2,
                 align=PP_ALIGN.CENTER)
        # Arrow from surface DOWN to handler
        arrow(s, x + Inches(1.1), y + Inches(1.4),
              cx + (x - cx + Inches(1.1) - cx) * 0.05,  # slight fan-in
              cy - handler_h / 2 - Inches(0.05),
              color=color, weight=1.5)

    # Three downstream emitters
    downstream = [
        ("Tool result · widgets · markdown", BAD, Inches(2.0), Inches(6.3)),
        ("MCP UI blocks (text/html)",        BAD, cx - Inches(1.7), Inches(6.3)),
        ("Same handler, same audit entry",   BAD, W - Inches(4.6), Inches(6.3)),
    ]
    for label, color, x, y in downstream:
        rounded(s, x, y, Inches(2.6), Inches(0.5), BRAND_TINT, radius=0.04)
        text_box(s, x, y, Inches(2.6), Inches(0.5), label, size=10,
                 color=BRAND_DEEP, bold=True, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        # Arrow from handler down to each
        arrow(s, cx + (x + Inches(1.3) - cx) * 0.05,
              cy + handler_h / 2 + Inches(0.05),
              x + Inches(1.3), y, color=color, weight=1.2, dashed=True)

    add_speaker_notes(s,
        "This slide is the architectural payoff. Three surfaces from one handler — that's why ToolDef is framework-agnostic data.",
        "Demo the eDiscovery review tool: same code path, three viewers see the same custodian card.")
    return s


def slide_lib_registries_v2(prs):
    """NEW: hub-and-spoke registry visualisation, replacing the 3-column slide."""
    s = add(prs)
    slide_chrome(s, "5 · THE LIBRARY", "13 registries — hub & spokes", "Architects")

    text_box(s, Inches(0.7), Inches(1.65), W - Inches(1.4), Inches(0.4),
             "All 13 implement one Registry<TDef> interface — same MFE-aware teardown, same signal contract, same conformance test surface.",
             size=12, color=TEXT_2)

    # Central hub
    cx, cy = W / 2, Inches(4.5)
    hub_r = Inches(1.0)
    rounded(s, cx - hub_r, cy - hub_r * 0.5, hub_r * 2, hub_r,
            BRAND_DEEP, radius=0.5)
    text_box(s, cx - hub_r, cy - hub_r * 0.5, hub_r * 2, hub_r,
             "Registry<TDef>", size=12, bold=True, color=SURFACE_2,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 13 spokes laid out in two arcs (top + bottom)
    import math
    nodes = [
        # Tier 1 - core (M1-M3)
        ("Tool",       BRAND, "core"),
        ("Component",  BRAND, "core"),
        ("Capability", BRAND, "core"),
        ("Backend",    BRAND, "core"),
        ("Mfe",        BRAND, "core"),
        # Tier 2 - extended (M4-M5)
        ("Action",     OK, "extended"),
        ("Intent",     OK, "extended"),
        ("Form",       OK, "extended"),
        ("DataSource", OK, "extended"),
        # Tier 3 - seams (M4-M5)
        ("Validation", WARN, "seam"),
        ("Persistence",WARN, "seam"),
        ("Layout",     WARN, "seam"),
        ("Schema",     WARN, "seam"),
    ]
    # Position 5 above and 8 below (tier 2+3) to keep tiers grouped
    above = nodes[:5]
    below = nodes[5:]
    radius_x = Inches(4.2); radius_y = Inches(2.1)

    # Above arc (top half) — spread across angles 200-340 (left-up around to right-up)
    for i, (name, color, _) in enumerate(above):
        # angle in degrees, sweep from 220 to 320 over 5 nodes
        a = math.radians(220 + i * (100 / max(1, len(above) - 1)))
        nx = cx + radius_x * math.cos(a) / Inches(1) * Inches(1)
        ny = cy + radius_y * math.sin(a) / Inches(1) * Inches(1)
        chip_node(s, nx - Inches(0.6), ny - Inches(0.18), Inches(1.2),
                  Inches(0.36), name, fill=SURFACE, fg=color, border=color, size=11)
        arrow(s, cx, cy, nx - Inches(0.6) + Inches(1.2) / 2 if nx < cx else nx - Inches(0.6),
              ny, color=color, weight=1)

    # Below arc — angles 20-160 spread across 8 nodes
    for i, (name, color, _) in enumerate(below):
        a = math.radians(20 + i * (140 / max(1, len(below) - 1)))
        nx = cx + radius_x * math.cos(a) / Inches(1) * Inches(1)
        ny = cy + radius_y * math.sin(a) / Inches(1) * Inches(1)
        chip_node(s, nx - Inches(0.6), ny - Inches(0.18), Inches(1.2),
                  Inches(0.36), name, fill=SURFACE, fg=color, border=color, size=11)
        arrow(s, cx, cy + Inches(0.05), nx - Inches(0.6) + Inches(0.6),
              ny, color=color, weight=1)

    # Tier legend
    leg_y = Inches(7.0)
    legend = [
        ("Core (M1–M3)", BRAND, Inches(0.7)),
        ("Extended (M4–M5)", OK, Inches(4.5)),
        ("Seams", WARN, Inches(8.3)),
    ]
    for label, color, x in legend:
        rounded(s, x, leg_y - Inches(0.12), Inches(0.18), Inches(0.18),
                color, radius=0.5)
        text_box(s, x + Inches(0.25), leg_y - Inches(0.18), Inches(3),
                 Inches(0.3), label, size=11, bold=True, color=color)

    add_speaker_notes(s,
        "Hub-and-spoke shows the architectural unity — the base class is the hub, every registry is a spoke with the same shape.",
        "Tier 1 is required for any agentic UI; tiers 2 and 3 are opt-in.")
    return s


def slide_agui_event_timeline(prs):
    """NEW: AG-UI event timeline — replaces the bullet-list of events."""
    s = add(prs)
    slide_chrome(s, "2 · AG-UI", "Event timeline — one chat turn", "Developers")

    text_box(s, Inches(0.7), Inches(1.65), W - Inches(1.4), Inches(0.4),
             "Time flows left → right. Pills are the events on the wire; vertical position groups them by class.",
             size=11, color=MUTED)

    # Time axis
    axis_y = Inches(6.4)
    arrow(s, Inches(0.9), axis_y, W - Inches(0.9), axis_y,
          color=BORDER, weight=1.2)
    for tick_x, tick_label in [
        (Inches(1.5), "0 ms"), (Inches(4.0), "1.2 s"),
        (Inches(8.0), "3.4 s"), (Inches(12.0), "5.1 s"),
    ]:
        rect(s, tick_x, axis_y - Inches(0.05), Inches(0.02), Inches(0.1), MUTED)
        text_box(s, tick_x - Inches(0.5), axis_y + Inches(0.1), Inches(1),
                 Inches(0.3), tick_label, size=9, color=MUTED,
                 align=PP_ALIGN.CENTER)

    # Lifecycle row (top)
    lc_y = Inches(2.3)
    text_box(s, Inches(0.7), lc_y - Inches(0.05), Inches(1.6),
             Inches(0.3), "LIFECYCLE", size=9, bold=True, color=BRAND)
    event_pill(s, Inches(1.5),  lc_y, Inches(1.5), "RUN_STARTED", color=BRAND)
    event_pill(s, Inches(11.4), lc_y, Inches(1.5), "RUN_FINISHED", color=BRAND)

    # Text row
    tx_y = Inches(3.0)
    text_box(s, Inches(0.7), tx_y - Inches(0.05), Inches(1.6),
             Inches(0.3), "TEXT", size=9, bold=True, color=INFO)
    event_pill(s, Inches(2.0),  tx_y, Inches(1.4), "MESSAGE_START", color=INFO)
    for i, x in enumerate([Inches(2.5 + i * 0.4) for i in range(4)]):
        event_pill(s, x + Inches(1.1), tx_y, Inches(0.4), "Δ", color=INFO)
    event_pill(s, Inches(5.5),  tx_y, Inches(1.4), "MESSAGE_END",   color=INFO)

    # Tool-call row
    tc_y = Inches(3.7)
    text_box(s, Inches(0.7), tc_y - Inches(0.05), Inches(1.6),
             Inches(0.3), "TOOL CALL", size=9, bold=True, color=OK)
    event_pill(s, Inches(7.0), tc_y, Inches(1.5), "CALL_START", color=OK)
    event_pill(s, Inches(8.6), tc_y, Inches(1.0), "ARGS",       color=OK)
    event_pill(s, Inches(9.7), tc_y, Inches(1.4), "CALL_END",   color=OK)
    event_pill(s, Inches(11.2),tc_y, Inches(1.5), "RESULT",     color=OK)

    # Generative row
    gu_y = Inches(4.4)
    text_box(s, Inches(0.7), gu_y - Inches(0.05), Inches(1.6),
             Inches(0.3), "GENERATIVE UI", size=9, bold=True, color=WARN)
    event_pill(s, Inches(11.1), gu_y, Inches(1.7), "WIDGET_RENDER", color=WARN)

    # Post-tool text
    pt_y = Inches(5.1)
    text_box(s, Inches(0.7), pt_y - Inches(0.05), Inches(1.6),
             Inches(0.3), "POST-TOOL TEXT", size=9, bold=True, color=INFO)
    event_pill(s, Inches(11.7), pt_y, Inches(1.2), "MESSAGE", color=INFO)

    # Vertical guides at tick positions
    for x in (Inches(1.5), Inches(4.0), Inches(8.0), Inches(12.0)):
        arrow(s, x, lc_y - Inches(0.1), x, axis_y,
              color=BORDER, weight=0.5, dashed=True, head_w=0, head_l=0)

    # Caption
    text_box(s, Inches(0.7), Inches(7.0), W - Inches(1.4), Inches(0.4),
             "Adapter @maverick/agentic-ui/ag-ui maps each event 1:1 onto the AgenticEvent union.",
             size=11, color=MUTED, align=PP_ALIGN.CENTER, bold=True)
    add_speaker_notes(s,
        "Real-shape timeline — every event a stable record. This is what makes AG-UI audit-friendly.",
        "Note widget_render is synthesised from show-component tool-call result; it doesn't exist on the wire.")
    return s


def slide_lib_registries(prs):
    s = add(prs)
    slide_chrome(s, "5 · THE LIBRARY", "13 registries — the data model", "Architects")

    text_box(s, Inches(0.7), Inches(1.7), W - Inches(1.4), Inches(0.4),
             "All 13 implement one Registry<TDef> interface — same MFE-aware teardown, same signal contract, same conformance test surface.",
             size=12, color=TEXT_2)

    # Tier columns
    col_w = (W - Inches(1.4) - Inches(0.4)) / 3
    titles = ["Core (M1–M3)\nchat shell depends on these",
              "Extended (M4–M5)\nfull agent-driven UI",
              "Seams (M4–M5)\ninterface + thin default"]
    bodies = [
        ["ToolRegistry",
         "ComponentRegistry  (alias WidgetRegistry)",
         "CapabilityRegistry",
         "BackendRegistry",
         "MfeRegistry  (external)"],
        ["ActionRegistry",
         "IntentRegistry",
         "FormRegistry",
         "DataSourceRegistry"],
        ["ValidationRegistry",
         "PersistenceRegistry",
         "LayoutRegistry",
         "SchemaTransformerRegistry"],
    ]
    colors = [BRAND, OK, WARN]
    for i in range(3):
        x = Inches(0.7) + i * (col_w + Inches(0.2))
        rounded(s, x, Inches(2.3), col_w, Inches(4.5), SURFACE, line=BORDER, radius=0.03)
        rect(s, x, Inches(2.3), col_w, Inches(0.06), colors[i])
        text_box(s, x + Inches(0.2), Inches(2.45), col_w - Inches(0.4),
                 Inches(0.7), titles[i], size=12, bold=True, color=colors[i])
        bullets(s, x + Inches(0.2), Inches(3.25), col_w - Inches(0.4),
                Inches(3.4), bodies[i], size=12, gap=6, marker="—",
                marker_color=colors[i])
    add_speaker_notes(s,
        "13 registries sound like a lot — pay attention to the tiers. Most apps will only ever populate Tool + Component.",
        "Adding a new registry is ~30 LOC because they all share a base; cost per registry is sublinear.")
    return s


def slide_lib_backend(prs):
    s = add(prs)
    slide_chrome(s, "5 · THE LIBRARY", "AgenticBackend — the abstraction", "Developers")

    code_block(s, Inches(0.7), Inches(1.7), W - Inches(1.4), Inches(3.7), [
        "export interface AgenticBackend {",
        "  readonly id: string;",
        "  readonly capabilities: {",
        "    streaming: boolean;",
        "    clientTools: boolean;",
        "    generativeUi: boolean;",
        "    uiActions: boolean;     // A2UI",
        "  };",
        "  run(input: AgenticRunInput): AsyncIterable<AgenticEvent>;",
        "  reset?(threadId: string): Promise<void>;",
        "}",
        "",
        "// 12 events: run-started, text-delta, tool-call-*, widget-render, ui-action, …",
        "export type AgenticEvent =",
        "  | { type: 'run-started'; threadId: string; runId: string }",
        "  | { type: 'text-delta'; messageId: string; delta: string }",
        "  | { type: 'tool-call-result'; toolCallId: string; result: unknown }",
        "  | { type: 'widget-render'; widgetCallId: string; name: string; props: unknown }",
        "  | { type: 'ui-action'; actionId: string; op: string; payload: unknown }",
        "  | … ;",
    ])

    text_box(s, Inches(0.7), Inches(5.7), W - Inches(1.4), Inches(1.2),
             "The chat shell sees ONLY this surface. AgUiBackend, HashbrownBackend, "
             "A2uiBackend implement it — capability flags drive feature-detection in the UI.",
             size=12, color=TEXT_2, align=PP_ALIGN.CENTER)
    add_speaker_notes(s,
        "Show this slide to engineers. The chat shell is protocol-agnostic; only the backend adapter knows the wire format.",
        "Adding a new protocol is implementing this one interface.")
    return s


def slide_lib_mfe(prs):
    s = add(prs)
    slide_chrome(s, "5 · THE LIBRARY", "Microfrontend federation — capability handoff", "Architects")

    bullets(s, Inches(0.7), Inches(1.7), W - Inches(1.4), Inches(2.5), [
        ("Discover.", "MfeRegistryClient.discover(env) — Spring Boot or static-JSON adapter."),
        ("Lazy-load.", "loadRemoteModule({ remoteName, exposedModule: './Capability' })."),
        ("Register.", "Remote's defineCapabilityModule pushes into ToolRegistry / ComponentRegistry / ActionRegistry."),
        ("Teardown.", "removeBySource('remote:bookings') runs across every registry in one pass; signals notify subscribers."),
    ], size=13, gap=6)

    rounded(s, Inches(0.7), Inches(4.3), Inches(5.85), Inches(2.6),
            BRAND_TINT, line=BORDER, radius=0.04)
    text_box(s, Inches(0.95), Inches(4.45), Inches(5.4), Inches(0.4),
             "Native Federation (default)", size=13, bold=True, color=BRAND_DEEP)
    bullets(s, Inches(0.95), Inches(4.85), Inches(5.4), Inches(2.0), [
        "Backed by @angular-architects/native-federation",
        "esbuild-compatible — Angular CLI default",
        "Smaller bundle, faster cold start"
    ], size=12, gap=4, marker="—", marker_color=BRAND_DEEP)

    rounded(s, Inches(6.85), Inches(4.3), Inches(5.85), Inches(2.6),
            RGBColor(0xFE, 0xF3, 0xC7), line=BORDER, radius=0.04)
    text_box(s, Inches(7.1), Inches(4.45), Inches(5.4), Inches(0.4),
             "Module Federation (peer)", size=13, bold=True, color=WARN)
    bullets(s, Inches(7.1), Inches(4.85), Inches(5.4), Inches(2.0), [
        "Backed by @module-federation/runtime",
        "webpack — for teams already on it",
        "Same CapabilityModule format · same capabilities.json"
    ], size=12, gap=4, marker="—", marker_color=WARN)
    add_speaker_notes(s,
        "Both federation paths are first-class — neither is a side branch.",
        "Switch via `ng add @maverick/agentic-ui --federation=native` vs `--federation=module-federation`.")
    return s


def slide_lib_mcp(prs):
    s = add(prs)
    slide_chrome(s, "5 · THE LIBRARY", "MCP integration — analyst workstation reach")

    text_box(s, Inches(0.7), Inches(1.7), W - Inches(1.4), Inches(0.5),
             "@maverick/agentic-ui-mcp — the same ToolDefs power Claude Desktop, Cursor, Zed.",
             size=15, color=TEXT_2)

    rounded(s, Inches(0.7), Inches(2.5), W - Inches(1.4), Inches(3.4),
            SURFACE_2, line=BORDER, radius=0.02)
    text_box(s, Inches(0.95), Inches(2.65), W - Inches(1.8), Inches(0.4),
             "One handler, multiple surfaces", size=14, bold=True, color=BRAND_DEEP)
    bullets(s, Inches(0.95), Inches(3.1), W - Inches(1.8), Inches(2.6), [
        "agenticTool({ ... }) — same factory the chat shell uses",
        "mcpToolBridge(tools) — exposes them as an MCP server with a Node binary",
        "MCP UI: tool results carry html field → text/html;profile=mcp-app blocks",
        "Per-user MCP server pattern — beforeCall stub auth → real OIDC swap-in",
        "Cookbook: paralegal privilege review in Claude Desktop using the eDiscovery toolset",
    ], size=12, gap=6)

    text_box(s, Inches(0.7), Inches(6.2), W - Inches(1.4), Inches(0.6),
             "Result: agents that work in your web app AND in your analysts' IDE — same audit trail, same governance.",
             size=12, color=BRAND_DEEP, bold=True, align=PP_ALIGN.CENTER)
    add_speaker_notes(s,
        "Big win for regulated industries — MCP brings agentic UI into Claude Desktop where many analysts already live.",
        "Demonstrate by walking through demo-mcp-server then opening Claude Desktop.")
    return s


def slide_lib_telemetry(prs):
    s = add(prs)
    slide_chrome(s, "5 · THE LIBRARY", "Observability — OTel from M1", "Architects")

    text_box(s, Inches(0.7), Inches(1.7), W - Inches(1.4), Inches(0.5),
             "Distributed tracing for chat shell → backend → SSE route → agent → LLM → tool execution.",
             size=14, color=TEXT_2)

    rows = [
        ("agentic.run",            "INTERNAL", "thread_id, run_id, backend.id"),
        ("agentic.backend.stream", "CLIENT",   "http.url, http.status_code, traceparent"),
        ("agentic.tool_call",      "INTERNAL", "tool.name, tool.source, success"),
        ("agentic.federation.load","INTERNAL", "remote_name, version, federation, load_ms"),
        ("agentic.registry.register","INTERNAL","registry.name, source"),
    ]
    rounded(s, Inches(0.7), Inches(2.5), W - Inches(1.4), Inches(3.7),
            SURFACE, line=BORDER, radius=0.02)
    rect(s, Inches(0.7), Inches(2.5), W - Inches(1.4), Inches(0.45), SURFACE_2)
    text_box(s, Inches(0.95), Inches(2.55), Inches(4.0), Inches(0.4),
             "Span", size=11, bold=True, color=MUTED)
    text_box(s, Inches(5.0), Inches(2.55), Inches(2.0), Inches(0.4),
             "Kind", size=11, bold=True, color=MUTED)
    text_box(s, Inches(7.5), Inches(2.55), Inches(5.0), Inches(0.4),
             "Key attributes", size=11, bold=True, color=MUTED)
    y = Inches(2.95)
    for span, kind, attrs in rows:
        text_box(s, Inches(0.95), y, Inches(4.0), Inches(0.5),
                 span, size=12, bold=True, color=BRAND_DEEP, font="Menlo")
        text_box(s, Inches(5.0), y, Inches(2.0), Inches(0.5),
                 kind, size=11, color=MUTED)
        text_box(s, Inches(7.5), y, Inches(5.0), Inches(0.5),
                 attrs, size=11, color=TEXT_2, font="Menlo")
        y += Inches(0.55)

    text_box(s, Inches(0.7), Inches(6.4), W - Inches(1.4), Inches(0.5),
             "W3C traceparent propagated across the SSE boundary — one trace covers UI to LLM.",
             size=11, color=MUTED, align=PP_ALIGN.CENTER, bold=True)
    add_speaker_notes(s,
        "Default sink is no-op; consumer opts in to OTel via @maverick/agentic-ui/otel.",
        "Bundle impact: zero unless they import the otel entry — measured in CI.")
    return s


def slide_lib_schematics(prs):
    s = add(prs)
    slide_chrome(s, "5 · THE LIBRARY", "Schematics — ng add → fully wired app", "Developers")

    text_box(s, Inches(0.7), Inches(1.7), W - Inches(1.4), Inches(0.4),
             "@maverick/agentic-ui-schematics — six generators ship with the package.",
             size=14, color=TEXT_2)

    code_block(s, Inches(0.7), Inches(2.2), W - Inches(1.4), Inches(2.5), [
        "$ ng add @maverick/agentic-ui \\",
        "    --backend=ag-ui \\",
        "    --mfe=host --federation=native \\",
        "    --registry=spring-boot \\",
        "    --server=mastra \\",
        "    --telemetry=otel",
        "",
        "Patches app.config.ts · seeds tools/widgets · scaffolds federation · ",
        "wires Mastra agent server · adds OTel collector config",
    ])

    text_box(s, Inches(0.7), Inches(5.0), W - Inches(1.4), Inches(0.4),
             "Per-artefact generators", size=14, bold=True, color=BRAND)
    bullets(s, Inches(0.7), Inches(5.5), W - Inches(1.4), Inches(1.7), [
        ("ng g …:tool",            "*.tool.ts with Zod schema scaffold; auto-registers"),
        ("ng g …:widget",          "Standalone component + Zod props schema"),
        ("ng g …:backend",         "Custom AgenticBackend with FakeAgenticBackend test"),
        ("ng g …:mfe-capability",  "Federation expose + capabilities.json manifest"),
        ("ng g …:agent-server",    "Mastra agent + AG-UI SSE route + memory store"),
        ("ng g …:chat-shell",      "Routed component using <mvk-chat-shell>"),
    ], size=12, gap=4)
    add_speaker_notes(s,
        "Show this slide to developers — the schematics are the value-multiplier.",
        "ng add takes a fresh Angular app to running chat in <2 minutes.")
    return s


# ── Examples (the demos in this repo) ───────────────────────────────────────

def slide_examples_overview(prs):
    s = add(prs)
    slide_chrome(s, "6 · EXAMPLES", "Five demos in this repo", "Developers")

    rows = [
        ("demo-monolith",        "Single-app AG-UI parity",                "M1 baseline · same shape as flights42"),
        ("demo-feature-tour",    "All 13 registries in one app",          "ActionRegistry, FormRegistry, IntentRegistry, DataSource"),
        ("demo-shell + remotes", "Native Federation host + bookings/loyalty/support",
                                 "Capability handoff · prefetchCapabilities"),
        ("demo-multi-agent",     "OrchestratorAgent + specialists",       "Sticky routing · ThreadStateStore"),
        ("demo-mcp-server",      "Tools as MCP server",                   "Claude Desktop · MCP UI · per-user instances"),
        ("demo-ediscovery",      "Enterprise reference app",              "Federated review remote · audit trail · 5 personas"),
    ]
    rounded(s, Inches(0.7), Inches(1.7), W - Inches(1.4), Inches(5.4),
            SURFACE, line=BORDER, radius=0.02)
    rect(s, Inches(0.7), Inches(1.7), W - Inches(1.4), Inches(0.5), SURFACE_2)
    text_box(s, Inches(0.95), Inches(1.78), Inches(4.0), Inches(0.4),
             "Demo", size=11, bold=True, color=MUTED)
    text_box(s, Inches(5.0), Inches(1.78), Inches(3.5), Inches(0.4),
             "What it shows", size=11, bold=True, color=MUTED)
    text_box(s, Inches(8.7), Inches(1.78), Inches(4.0), Inches(0.4),
             "Library features exercised", size=11, bold=True, color=MUTED)
    y = Inches(2.25)
    for name, what, feats in rows:
        text_box(s, Inches(0.95), y, Inches(4.0), Inches(0.6),
                 name, size=12, bold=True, color=BRAND_DEEP, font="Menlo")
        text_box(s, Inches(5.0), y, Inches(3.6), Inches(0.6),
                 what, size=11, color=TEXT_2)
        text_box(s, Inches(8.7), y, Inches(4.1), Inches(0.6),
                 feats, size=10, color=MUTED)
        y += Inches(0.78)
    add_speaker_notes(s,
        "Six demos. Pick one for each audience.",
        "demo-ediscovery is the headline — full enterprise pattern with federation + audit + multi-persona.")
    return s


def slide_example_ediscovery(prs):
    """eDiscovery case-study slide (v3) — strategy-deck stat-card layout
    with restrained typography. Lead with the four headline metrics."""
    s = add(prs)
    slide_chrome(s, "06 · CASE STUDY", "An enterprise reference app",
                 "Senior executives")

    text_box(s, Inches(0.7), Inches(1.85), W - Inches(1.4), Inches(0.5),
             "A federated eDiscovery shell — the demo we use to validate every "
             "library feature against a recognisable, regulated, complex domain.",
             size=14, color=TEXT_2)

    # Five headline stat cards across the top
    sw = (W - Inches(1.4) - Inches(0.8)) / 5
    sy = Inches(2.6)
    sh = Inches(1.3)
    stats = [
        ("8",    "phases shipped",               BRAND),
        ("18",   "tools · 4 specialists",        BRAND),
        ("3 + 1","MFE remotes + MCP server",     BRAND),
        ("13",   "registries exercised",         BRAND),
        ("✓",    "tamper-evident audit",         ACCENT),
    ]
    for i, (val, cap, col) in enumerate(stats):
        sx = Inches(0.7) + i * (sw + Inches(0.2))
        stat_card(s, sx, sy, sw, sh, val, cap, color=col)

    # Body — two columns: scope + library seams
    bx = Inches(0.7); by = Inches(4.2); bw = (W - Inches(1.4) - Inches(0.4)) / 2

    text_box(s, bx, by, bw, Inches(0.32),
             "WHAT THE DEMO SHIPS", size=10, bold=True, color=ACCENT)
    rect(s, bx, by + Inches(0.36), bw, Inches(0.012), BORDER)
    bullets(s, bx, by + Inches(0.5), bw, Inches(2.5), [
        "Custodian intake · legal-hold lifecycle",
        "Document review with privilege workflow",
        "Production assembly: create → bates → export",
        "Search + TAR classifier across the matter",
        "Tamper-evident chain-of-custody report",
        "MCP server for Claude Desktop / Cursor / Zed",
        "Persona-scoped tool surface (5 roles)",
    ], size=11, gap=4, marker_color=BRAND)

    rx = bx + bw + Inches(0.4)
    text_box(s, rx, by, bw, Inches(0.32),
             "LIBRARY SEAMS PUT TO WORK", size=10, bold=True, color=ACCENT)
    rect(s, rx, by + Inches(0.36), bw, Inches(0.012), BORDER)
    bullets(s, rx, by + Inches(0.5), bw, Inches(2.5), [
        "ToolRegistry · ComponentRegistry · CapabilityRegistry",
        "ActionRegistry — click-through navigation from widgets",
        "FormRegistry — productionConfigForm",
        "ValidationRegistry — Bates pattern conformance",
        "DataSourceRegistry — pluggable documentIndex",
        "Telemetry — every tool call audit-logged",
    ], size=12, gap=6, marker_color=BRAND)

    add_speaker_notes(s,
        "The eDiscovery reference is the headline asset.",
        "Open with the four stat cards — those are the talking points executives remember.")
    return s


# ── Decision matrix + benefits ──────────────────────────────────────────────

def slide_when_to_use_what(prs):
    s = add(prs)
    slide_chrome(s, "7 · DECISIONS", "When to use what — a guide")

    rows = [
        ("Single-app · regulated · audit-heavy",            "AG-UI",   BRAND),
        ("Multi-provider eval · rapid prototype",           "Hashbrown", OK),
        ("Agent navigates / fills forms / mutates store",   "A2UI (over AG-UI base)", WARN),
        ("Multiple teams · independent deploys",            "MFE federation",   INFO),
        ("Analyst desktops · IDE-native agents",            "MCP server",       BAD),
        ("Compliance + chain of custody",                   "Telemetry sink → audit log", BRAND_DEEP),
    ]
    rounded(s, Inches(0.7), Inches(1.7), W - Inches(1.4), Inches(5.4),
            SURFACE, line=BORDER, radius=0.02)
    rect(s, Inches(0.7), Inches(1.7), W - Inches(1.4), Inches(0.45), SURFACE_2)
    text_box(s, Inches(0.95), Inches(1.78), Inches(8.0), Inches(0.4),
             "Use case", size=11, bold=True, color=MUTED)
    text_box(s, Inches(9.0), Inches(1.78), Inches(3.5), Inches(0.4),
             "Reach for", size=11, bold=True, color=MUTED)
    y = Inches(2.3)
    for use, pick, color in rows:
        text_box(s, Inches(0.95), y + Inches(0.1), Inches(8.0), Inches(0.5),
                 use, size=13, color=TEXT)
        chip(s, Inches(9.0), y + Inches(0.15), pick,
             fg=SURFACE, bg=color, size=11, w=Inches(3.6))
        y += Inches(0.78)
    add_speaker_notes(s,
        "Don't try to pick one protocol. Pick by use case — the library lets you mix.")
    return s


def slide_benefits_execs(prs):
    s = add(prs)
    slide_chrome(s, "8 · BENEFITS", "For senior executives", "Senior executives")

    items = [
        ("Faster time-to-pilot.", "ng add → working chat in 2 minutes. 6 months of in-house plumbing avoided."),
        ("Vendor agnosticism.", "AG-UI today, Hashbrown tomorrow, A2UI when it ships v1 — same UI, swap the backend."),
        ("Regulated-domain ready.", "Audit-grade telemetry, persona scopes, MCP for analysts — built in, not bolted on."),
        ("Microfrontend native.", "Multiple teams contribute capabilities into one chat — federation is a first-class seam."),
        ("Open license.", "MIT. No SaaS lock-in. Internal forking encouraged."),
        ("Reference implementations.", "Six demos including a full eDiscovery app. Blueprint, not a brochure."),
    ]
    bullets(s, Inches(0.7), Inches(1.7), W - Inches(1.4), Inches(5.4),
            items, size=14, gap=12)
    add_speaker_notes(s, "Exec talk track. Six bullets, nothing technical — focus on derisking and time-to-value.")
    return s


def slide_benefits_architects(prs):
    s = add(prs)
    slide_chrome(s, "8 · BENEFITS", "For architects", "Architects")

    items = [
        ("Clean seams.", "AgenticBackend, MfeRegistrySource, Registry<TDef>, AgenticTelemetrySink — every contingency leans on one already in place."),
        ("Conformance-tested adapters.", "/testing entry runs the same suite against every backend; capability flags drive feature-detection."),
        ("Conflict policies + onDispose.", "Multiple teams may register the same tool name; namespace policies coexist; cleanup is automatic."),
        ("Observability by default.", "Distributed tracing across SSE; metrics; logs with trace_id correlation; opt-in OTel exporter."),
        ("Two federation paths.", "Native + Module Federation — no second-class citizen; same CapabilityModule shape; pick by bundler."),
        ("ADRs documented.", "Six ADRs cover backend abstraction, MCP, federation, conflict policies — context for future decisions."),
    ]
    bullets(s, Inches(0.7), Inches(1.7), W - Inches(1.4), Inches(5.4),
            items, size=13, gap=10)
    add_speaker_notes(s, "Architect framing. Highlight that every risk in §11 of the plan has a documented detection signal + contingency.")
    return s


def slide_benefits_devs(prs):
    s = add(prs)
    slide_chrome(s, "8 · BENEFITS", "For developers", "Developers")

    items = [
        ("Schematics.", "ng add bootstraps an entire workspace. ng g …:tool/widget/backend cuts hours of typing per artefact."),
        ("Signal-first APIs.", "ToolRegistry.signal(), CapabilityRegistry.signal() — reactive everywhere; OnPush by default."),
        ("Strict TypeScript.", "Zod-anchored payloads; tool args + props inferred; no any in the public API surface."),
        ("Testing harness.", "FakeAgenticBackend + harnessChat() in /testing; in-memory AgenticTelemetrySink with custom matchers."),
        ("Cookbook entries.", "Every milestone ships a how-to guide — federation at scale, MCP server, observability, persona scopes."),
        ("Live reload across MFEs.", "Both Native + Module Federation play nice with the dev server."),
    ]
    bullets(s, Inches(0.7), Inches(1.7), W - Inches(1.4), Inches(5.4),
            items, size=13, gap=10)
    add_speaker_notes(s, "Dev framing. Show the schematics demo at the end — typically the moment that lands the pitch.")
    return s


# ── Roadmap + CTA ───────────────────────────────────────────────────────────

def slide_roadmap(prs):
    """Library milestones (M1–M5) cross-referenced with the eDiscovery
    flagship phases that validated each one."""
    s = add(prs)
    slide_chrome(s, "10 · ROADMAP", "Library milestones · validated by eDiscovery",
                 "Architects")

    text_box(s, Inches(0.7), Inches(1.85), W - Inches(1.4), Inches(0.4),
             "Each library milestone has a corresponding eDiscovery phase that exercised it under enterprise load.",
             size=11, color=MUTED)

    rows = [
        ("M1", "AG-UI parity + ChatShell + /testing",
         "ToolRegistry, ComponentRegistry, BackendRegistry",
         "P0–P1", OK),
        ("M2", "Schematics + ng-add",
         "tool, widget, backend, agent-server, snapshot tests",
         "—",  OK),
        ("M3", "MFE federation + A2UI + MCP server adapter",
         "CapabilityRegistry, MfeRegistry, both federation paths",
         "P2, P6", OK),
        ("M4", "Extended registries + Validation + DataSource",
         "Action, Intent, Form, Validation, DataSource",
         "P2–P4", OK),
        ("M5", "Compliance + scope governance",
         "Tamper-evident audit chain, RegistryBase.setScopePolicy",
         "P5, P7–P8", OK),
    ]
    rounded(s, Inches(0.7), Inches(2.45), W - Inches(1.4), Inches(4.6),
            SURFACE, line=BORDER, radius=0.02)
    rect(s, Inches(0.7), Inches(2.45), W - Inches(1.4), Inches(0.45), SURFACE_2)
    text_box(s, Inches(0.95), Inches(2.53), Inches(0.6), Inches(0.4),
             "M",      size=10, bold=True, color=MUTED)
    text_box(s, Inches(1.7), Inches(2.53), Inches(3.5), Inches(0.4),
             "Theme",  size=10, bold=True, color=MUTED)
    text_box(s, Inches(5.5), Inches(2.53), Inches(5.0), Inches(0.4),
             "Ships",  size=10, bold=True, color=MUTED)
    text_box(s, Inches(11.0), Inches(2.53), Inches(1.4), Inches(0.4),
             "Validated by", size=10, bold=True, color=MUTED)
    y = Inches(3.05)
    for m, theme, ships, validated, color in rows:
        rounded(s, Inches(0.95), y + Inches(0.05), Inches(0.6), Inches(0.6),
                color, radius=0.5)
        text_box(s, Inches(0.95), y + Inches(0.05), Inches(0.6), Inches(0.6),
                 m, size=14, bold=True, color=SURFACE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, Inches(1.7), y + Inches(0.13), Inches(3.7), Inches(0.32),
                 theme,    size=12, bold=True, color=TEXT)
        text_box(s, Inches(1.7), y + Inches(0.46), Inches(3.7), Inches(0.32),
                 "",       size=9, color=MUTED)
        text_box(s, Inches(5.5), y + Inches(0.18), Inches(5.4), Inches(0.5),
                 ships,    size=11, color=TEXT_2)
        chip_node(s, Inches(11.0), y + Inches(0.18), Inches(1.5), Inches(0.34),
                  validated, fg=color,
                  fill=ACCENT_TINT if color == OK else SURFACE_2,
                  border=color, size=10)
        y += Inches(0.78)

    text_box(s, Inches(0.7), Inches(7.1), W - Inches(1.4), Inches(0.32),
             "All five milestones shipped. The eDiscovery flagship is the proof-shape under load.",
             size=11, bold=True, color=OK, align=PP_ALIGN.CENTER)

    add_speaker_notes(s,
        "Updated for v5 — every milestone is now shipped (was M1-M3 done / M4-M5 in flight).",
        "The 'validated by' column ties library work to the eDiscovery phase that put it under load.")
    return s


def slide_cta(prs):
    """CTA (v3) — high-contrast navy slab with restrained typography."""
    s = add(prs)
    rect(s, Inches(0), Inches(0), W, H, BRAND_DEEP)

    # Eyebrow
    text_box(s, Inches(0.7), Inches(0.6), Inches(8), Inches(0.32),
             "10 · CALL TO ACTION", size=11, bold=True, color=ACCENT)
    rect(s, Inches(0.7), Inches(1.0), Inches(1.0), Inches(0.025), ACCENT)

    # Title
    text_box(s, Inches(0.7), Inches(1.6), W - Inches(1.4), Inches(1.0),
             "Where you can plug in",
             size=42, bold=True, color=SURFACE)
    text_box(s, Inches(0.7), Inches(2.5), W - Inches(1.4), Inches(0.5),
             "Three concrete next steps — pick the row that matches your role.",
             size=14, color=RGBColor(0xCB, 0xD0, 0xD9))

    # Three rows — each as a thin row with hairline divider
    rows = [
        ("01",
         "EXECUTIVES",
         "Authorise a two-week pilot.",
         "Agentic chat retro-fit into your most-used internal app, MIT licensed."),
        ("02",
         "ARCHITECTS",
         "Read ADR-006 (MCP) + ADR-007 (eDiscovery). Spike a remote in your repo.",
         "Conformance suite + cookbook ship in /testing."),
        ("03",
         "DEVELOPERS",
         "ng add @maverick/agentic-ui. ng g …:tool. ng serve.",
         "30 minutes from clone to running chat — schematics scaffold the rest."),
    ]
    y = Inches(3.4)
    for num, who, action, why in rows:
        rect(s, Inches(0.7), y, W - Inches(1.4), Inches(0.012),
             RGBColor(0x35, 0x42, 0x5A))
        text_box(s, Inches(0.7), y + Inches(0.2), Inches(0.7), Inches(0.4),
                 num, size=24, bold=True, color=ACCENT)
        text_box(s, Inches(1.5), y + Inches(0.27), Inches(2.4), Inches(0.3),
                 who, size=10, bold=True, color=ACCENT)
        text_box(s, Inches(4.0), y + Inches(0.18), W - Inches(4.7), Inches(0.4),
                 action, size=15, bold=True, color=SURFACE)
        text_box(s, Inches(4.0), y + Inches(0.6), W - Inches(4.7), Inches(0.4),
                 why, size=11, color=RGBColor(0xCB, 0xD0, 0xD9))
        y += Inches(1.05)

    # Closing rule + caption
    rect(s, Inches(0.7), y + Inches(0.05), W - Inches(1.4), Inches(0.012),
         RGBColor(0x35, 0x42, 0x5A))
    text_box(s, Inches(0.7), H - Inches(0.6), W - Inches(1.4), Inches(0.3),
             "Repository · examples · cookbook · ADRs included     "
             "Questions? Reach the team.",
             size=10, color=RGBColor(0x8C, 0x96, 0xA8), align=PP_ALIGN.CENTER)
    add_speaker_notes(s,
        "Close. Three asks, one per audience.",
        "Tone matches the deck — restrained, navy, no decoration.")
    return s


def slide_resources(prs):
    s = add(prs)
    slide_chrome(s, "RESOURCES", "Where to go from here")

    text_box(s, Inches(0.7), Inches(1.7), Inches(5.5), Inches(0.4),
             "Repository layout", size=14, bold=True, color=BRAND_DEEP)
    code_block(s, Inches(0.7), Inches(2.15), Inches(5.7), Inches(4.6), [
        "ag_ui_maverick/",
        "├── projects/",
        "│   ├── agentic-ui/         # publishable lib",
        "│   ├── agentic-ui-server/  # Node companion",
        "│   └── agentic-ui-mcp/     # MCP adapter",
        "├── examples/",
        "│   ├── demo-monolith/",
        "│   ├── demo-shell/  (host)",
        "│   ├── demo-remote-bookings/",
        "│   ├── demo-multi-agent/",
        "│   ├── demo-mcp-server/",
        "│   └── demo-ediscovery-{shared,server,shell,review}/",
        "├── docs/",
        "│   ├── adr/                # decision records",
        "│   ├── cookbook/           # how-to guides",
        "│   └── architecture/       # system docs",
        "└── PLAN.md                 # full roadmap",
    ])

    text_box(s, Inches(6.7), Inches(1.7), Inches(6.0), Inches(0.4),
             "Recommended reading order", size=14, bold=True, color=BRAND_DEEP)
    items = [
        ("README.md",            "5-minute orientation"),
        ("PLAN.md",              "Full architecture + 5-tier registry model"),
        ("ADR-006",              "MCP server-side adapter design"),
        ("docs/cookbook/",       "How-to guides for federation, MCP, OTel"),
        ("examples/demo-feature-tour", "All 13 registries in one app"),
        ("examples/demo-ediscovery",   "Enterprise reference"),
        ("@maverick/agentic-ui/testing", "Conformance suite + harness"),
    ]
    bullets(s, Inches(6.7), Inches(2.15), Inches(6.0), Inches(4.6),
            items, size=12, gap=8)
    add_speaker_notes(s, "Last slide. Hand-off to the audience — every box on the left has its own README.")
    return s


# ── v5 — compliance, governance, MCP reach (Phases 5-8) ───────────────────

def slide_chain_of_custody(prs):
    """Phase 5 story — tamper-evident audit chain. The diagram shows
    three audit events linked by hash, with one "broken" event high-
    lighted to demonstrate what verification catches."""
    s = add(prs)
    slide_chrome(s, "06 · COMPLIANCE",
                 "Chain of custody — every action, every actor, hashed",
                 "Senior executives")

    text_box(s, Inches(0.7), Inches(1.85), W - Inches(1.4), Inches(0.5),
             "Phase 5. Each tool call writes one audit event; each event "
             "links to the previous via a chain hash. Tampering anywhere "
             "in the chain breaks verification.",
             size=14, color=TEXT_2)

    # Three event blocks linked by chain hashes
    blk_w = Inches(3.6)
    blk_h = Inches(2.0)
    blk_y = Inches(2.7)
    gap = Inches(0.3)
    start_x = Inches(0.7)

    blocks = [
        ("Event #1 — issued",
         "actor:  eleanor.vance@firm\n"
         "action: hold.issued\n"
         "target: HOLD-001\n"
         "ts:     2026-02-16T14:00Z",
         "0000…", "a3f8d2", OK),
        ("Event #2 — tagged",
         "actor:  paralegal-1@firm\n"
         "action: document.tag.added\n"
         "target: DOC-7891234\n"
         "ts:     2026-04-21T09:14Z",
         "a3f8d2", "b7e1f4", OK),
        ("Event #3 — exported",
         "actor:  associate-2@firm\n"
         "action: production.delivered\n"
         "target: PROD-A1B2\n"
         "ts:     2026-04-30T16:55Z",
         "b7e1f4", "c2d5a1", OK),
    ]
    for i, (title, body, prev, this_h, color) in enumerate(blocks):
        x = start_x + i * (blk_w + gap)
        rounded(s, x, blk_y, blk_w, blk_h, SURFACE, line=BORDER, radius=0.04)
        rect(s, x, blk_y, blk_w, Inches(0.05), color)
        text_box(s, x + Inches(0.2), blk_y + Inches(0.15),
                 blk_w - Inches(0.4), Inches(0.32),
                 title, size=12, bold=True, color=color)
        text_box(s, x + Inches(0.2), blk_y + Inches(0.55),
                 blk_w - Inches(0.4), Inches(0.95),
                 body, size=10, color=TEXT_2, font="Menlo")
        # Hash row at bottom
        rect(s, x + Inches(0.2), blk_y + blk_h - Inches(0.45),
             blk_w - Inches(0.4), Inches(0.012), DIVIDER)
        text_box(s, x + Inches(0.2), blk_y + blk_h - Inches(0.4),
                 Inches(1.4), Inches(0.3),
                 f"prev: {prev}", size=8, color=MUTED, font="Menlo")
        text_box(s, x + blk_w - Inches(1.6), blk_y + blk_h - Inches(0.4),
                 Inches(1.4), Inches(0.3),
                 f"chain: {this_h}", size=8, color=color, font="Menlo", bold=True)
        # Arrow to next block
        if i < 2:
            arrow(s, x + blk_w, blk_y + blk_h / 2,
                  x + blk_w + gap, blk_y + blk_h / 2,
                  color=ACCENT, weight=2)

    # Verification badge below
    badge_y = Inches(5.0)
    rounded(s, Inches(0.7), badge_y, W - Inches(1.4), Inches(0.85),
            ACCENT_TINT, line=BORDER, radius=0.04)
    rect(s, Inches(0.7), badge_y, Inches(0.06), Inches(0.85), OK)
    text_box(s, Inches(0.95), badge_y + Inches(0.13),
             W - Inches(1.7), Inches(0.32),
             "✓  CHAIN VERIFIED — 3 of 3 hashes recompute · head c2d5a1",
             size=14, bold=True, color=OK)
    text_box(s, Inches(0.95), badge_y + Inches(0.45),
             W - Inches(1.7), Inches(0.32),
             "If any event were edited or deleted, the next event's `prev` link would fail to recompute "
             "and the integrity badge on the Audit Trail page would flip to broken.",
             size=10, color=TEXT_2)

    # Why exec band
    why_y = Inches(6.1)
    text_box(s, Inches(0.7), why_y, Inches(2.5), Inches(0.32),
             "WHY THIS MATTERS", size=10, bold=True, color=ACCENT)
    rect(s, Inches(0.7), why_y + Inches(0.36), Inches(11.9),
         Inches(0.012), BORDER)
    text_box(s, Inches(0.7), why_y + Inches(0.5), W - Inches(1.4), Inches(0.45),
             "Defensibility is the unit of trust opposing counsel and regulators measure. "
             "The chain is built into every tool call — Claude Desktop and the web app both "
             "contribute to the same chain. No second system to reconcile.",
             size=11, color=TEXT_2)

    add_speaker_notes(s,
        "Hash function is FNV-1a in the demo (deterministic for testing); SHA-256 + HSM-backed HMAC is the prod swap.",
        "Phase 5 ships 1 tool (generateChainOfCustodyReport), 1 widget, the chainHash field, and the integrity badge on the Audit Trail page.",
        "The defensibility story lands here for execs — chain integrity = a deliverable for opposing counsel.")
    return s


def slide_mcp_reach(prs):
    """Phase 6 story — MCP reach into analyst workstations.
    Diagram: same ToolDef literals → web app + Claude Desktop, with a
    single audit log under both surfaces."""
    s = add(prs)
    slide_chrome(s, "07 · MCP REACH",
                 "One toolset. Three surfaces. One audit log.",
                 "Architects")

    text_box(s, Inches(0.7), Inches(1.85), W - Inches(1.4), Inches(0.5),
             "Phase 6. Paralegals run privilege review from Claude Desktop / Cursor / "
             "Zed without opening the web app — same handlers, same audit chain (Phase 5).",
             size=14, color=TEXT_2)

    # Center: the shared ToolDef
    cx, cy = W / 2, Inches(4.4)
    th, tw = Inches(1.1), Inches(3.4)
    rounded(s, cx - tw / 2, cy - th / 2, tw, th, BRAND, radius=0.06)
    text_box(s, cx - tw / 2, cy - th / 2, tw, Inches(0.35),
             "SHARED", size=10, bold=True, color=BRAND_TINT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, cx - tw / 2, cy - th / 2 + Inches(0.3), tw, Inches(0.5),
             "ToolDef[]",
             size=22, bold=True, color=SURFACE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font="Menlo")
    text_box(s, cx - tw / 2, cy - th / 2 + Inches(0.75), tw, Inches(0.3),
             "5 review tools · framework-agnostic literals",
             size=10, color=BRAND_TINT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Top three surfaces
    surf_y = Inches(2.65)
    surfaces = [
        ("Web app", "<mvk-chat-shell> @ :4300\nNgComponentOutlet widgets",
         BRAND_DEEP, Inches(0.9)),
        ("Claude Desktop", "Stdio MCP server\ntext/html;profile=mcp-app",
         ACCENT, cx - Inches(2.0)),
        ("Cursor / Zed", "Same MCP server\nMCP-UI HTML cards",
         INFO, W - Inches(4.9)),
    ]
    for label, desc, color, x in surfaces:
        rounded(s, x, surf_y, Inches(4.0), Inches(1.0), SURFACE,
                line=color, radius=0.04)
        rect(s, x, surf_y, Inches(4.0), Inches(0.06), color)
        text_box(s, x + Inches(0.2), surf_y + Inches(0.15),
                 Inches(3.8), Inches(0.32),
                 label, size=13, bold=True, color=color)
        text_box(s, x + Inches(0.2), surf_y + Inches(0.5),
                 Inches(3.8), Inches(0.55),
                 desc, size=10, color=TEXT_2)
        # Arrow down to ToolDef
        arrow(s, x + Inches(2.0), surf_y + Inches(1.0),
              cx + (x + Inches(2.0) - cx) * 0.15, cy - th / 2 - Inches(0.05),
              color=color, weight=1.4)

    # Bottom: shared audit chain
    audit_y = Inches(5.4)
    rounded(s, Inches(0.7), audit_y, W - Inches(1.4), Inches(1.1),
            ACCENT_TINT, line=BORDER, radius=0.04)
    rect(s, Inches(0.7), audit_y, Inches(0.06), Inches(1.1), ACCENT)
    text_box(s, Inches(0.95), audit_y + Inches(0.15),
             W - Inches(1.7), Inches(0.35),
             "ONE AUDIT CHAIN — Phase 5 hash links every event regardless of surface",
             size=12, bold=True, color=ACCENT)
    text_box(s, Inches(0.95), audit_y + Inches(0.55),
             W - Inches(1.7), Inches(0.5),
             "appendAudit(...) auto-stamps chainHash + prevHash. The chain-of-custody "
             "report covers Claude-Desktop-driven mutations and web-app-driven mutations as one.",
             size=10, color=TEXT_2)
    # Arrow from ToolDef DOWN to audit
    arrow(s, cx, cy + th / 2, cx, audit_y, color=ACCENT,
          weight=1.5, dashed=True)

    # Why band
    why_y = Inches(6.7)
    text_box(s, Inches(0.7), why_y, Inches(11), Inches(0.32),
             "WHERE IT MATTERS — analysts who already live in Claude Desktop don't switch contexts to do privilege review.",
             size=11, bold=True, color=ACCENT)

    add_speaker_notes(s,
        "ToolDef literals are framework-agnostic — no Angular at MCP runtime.",
        "Per-user attribution via env vars (MVK_USER, MVK_MATTER) is the cookbook pattern for spawning one MCP server per user.",
        "Show this slide right after the chain-of-custody slide — the punchline lands when execs see two surfaces feeding one chain.")
    return s


def slide_persona_scope_policy(prs):
    """Phase 7+8 story — persona scopes from consumer-side shim
    (P7) → first-class library API (P8). Shows the simplification."""
    s = add(prs)
    slide_chrome(s, "08 · GOVERNANCE",
                 "Persona scopes — from consumer shim to library primitive",
                 "Architects")

    text_box(s, Inches(0.7), Inches(1.85), W - Inches(1.4), Inches(0.5),
             "Phases 7 → 8. The role allow-list moved from chat-shell-only "
             "filter to a registry-base policy that filters every read.",
             size=14, color=TEXT_2)

    # Two-column: Phase 7 vs Phase 8
    col_w = (W - Inches(1.4) - Inches(0.4)) / 2

    # Phase 7
    p7_x = Inches(0.7)
    p7_y = Inches(2.55)
    rounded(s, p7_x, p7_y, col_w, Inches(3.4), SURFACE, line=BORDER, radius=0.04)
    rect(s, p7_x, p7_y, col_w, Inches(0.06), MUTED)
    text_box(s, p7_x + Inches(0.2), p7_y + Inches(0.18),
             col_w - Inches(0.4), Inches(0.32),
             "PHASE 7 · CONSUMER-SIDE SHIM", size=11, bold=True, color=MUTED)
    text_box(s, p7_x + Inches(0.2), p7_y + Inches(0.55),
             col_w - Inches(0.4), Inches(0.4),
             "TOOL_FILTER + useFactory composition",
             size=12, bold=True, color=TEXT)
    code_block(s, p7_x + Inches(0.2), p7_y + Inches(1.0),
               col_w - Inches(0.4), Inches(2.2), [
        "{ provide: TOOL_FILTER,",
        "  useFactory: () =>",
        "    personaToolFilter(",
        "      inject(PersonaService),",
        "      keywordToolFilter({",
        "        maxTools: 12,",
        "        floor: 5,",
        "      }),",
        "    ),",
        "}",
    ])

    # Phase 8
    p8_x = p7_x + col_w + Inches(0.4)
    rounded(s, p8_x, p7_y, col_w, Inches(3.4), SURFACE, line=BORDER, radius=0.04)
    rect(s, p8_x, p7_y, col_w, Inches(0.06), OK)
    text_box(s, p8_x + Inches(0.2), p7_y + Inches(0.18),
             col_w - Inches(0.4), Inches(0.32),
             "PHASE 8 · LIBRARY API", size=11, bold=True, color=OK)
    text_box(s, p8_x + Inches(0.2), p7_y + Inches(0.55),
             col_w - Inches(0.4), Inches(0.4),
             "RegistryBase.setScopePolicy",
             size=12, bold=True, color=TEXT)
    code_block(s, p8_x + Inches(0.2), p7_y + Inches(1.0),
               col_w - Inches(0.4), Inches(2.2), [
        "// app.config.ts",
        "provideAppInitializer(() => {",
        "  const persona = inject(PersonaService);",
        "  inject(ToolRegistry).setScopePolicy(",
        "    (entry) =>",
        "      persona.canInvoke(",
        "        persona.active(),",
        "        entry.name,",
        "      ),",
        "  );",
        "});",
    ])

    # Bottom row: stat cards
    sc_y = Inches(6.15)
    sc_w = (W - Inches(1.4) - Inches(0.6)) / 4
    stats = [
        ("13", "registries gain scope filtering", BRAND),
        ("~30", "lines of consumer code deleted", OK),
        ("1", "policy installs once at boot", BRAND),
        ("✓", "every reader sees the same view", ACCENT),
    ]
    for i, (val, cap, col) in enumerate(stats):
        sx = Inches(0.7) + i * (sc_w + Inches(0.2))
        stat_card(s, sx, sc_y, sc_w, Inches(1.05), val, cap, color=col)

    add_speaker_notes(s,
        "Phase 7 shipped a chat-shell-only filter. Phase 8 lifts it to RegistryBase so 13 registries share the seam.",
        "The diff was about -30 LOC in the consumer for a one-line setScopePolicy call.",
        "ADR-008 documents the trade-offs: filter on read (not register) so federated remotes stay portable.")
    return s


def slide_compliance_summary(prs):
    """v5 — single-pane summary of Phases 5+6+7+8 for execs who skip the diagrams."""
    s = add(prs)
    slide_chrome(s, "09 · ENTERPRISE READINESS",
                 "Four phases, four governance primitives",
                 "Senior executives")

    text_box(s, Inches(0.7), Inches(1.85), W - Inches(1.4), Inches(0.4),
             "Each row pairs a regulated-domain need with the library primitive that addresses it — already shipped.",
             size=12, color=TEXT_2)

    rows = [
        ("Tamper-evident audit",
         "Chain hash on every event · live integrity verifier",
         "Phase 5 · `appendAudit` + `verifyAuditChain`", OK),
        ("Analyst-desktop reach",
         "Same toolset in Claude Desktop / Cursor · same audit log",
         "Phase 6 · `@maverick/agentic-ui-mcp`", ACCENT),
        ("Persona-scoped tool surface",
         "Hide tools the role can't invoke before the LLM ever sees them",
         "Phase 7 · consumer-side filter, then…", BRAND),
        ("Registry-wide scope policy",
         "13 registries (Tool / Action / Form / DataSource / …) share one filter",
         "Phase 8 · `RegistryBase.setScopePolicy(policy)`", BRAND_DEEP),
    ]

    rh = Inches(0.95)
    y = Inches(2.55)
    for need, what, ref, color in rows:
        rounded(s, Inches(0.7), y, W - Inches(1.4), rh - Inches(0.08),
                SURFACE, line=BORDER, radius=0.04)
        rect(s, Inches(0.7), y, Inches(0.06), rh - Inches(0.08), color)
        text_box(s, Inches(0.95), y + Inches(0.13),
                 Inches(3.8), Inches(0.4),
                 need, size=14, bold=True, color=color)
        text_box(s, Inches(0.95), y + Inches(0.5),
                 Inches(3.8), Inches(0.32),
                 ref, size=9, color=MUTED, font="Menlo")
        text_box(s, Inches(4.9), y + Inches(0.27),
                 W - Inches(5.6), Inches(0.5),
                 what, size=12, color=TEXT_2)
        y += rh

    add_speaker_notes(s,
        "Closing the compliance section — read top-to-bottom, each row stands on its own.",
        "Use this slide if you're tight on time and skipping the chain-of-custody / MCP / scope-policy graphics.")
    return s


# ── Compose ─────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # Each entry: (builder, section-label, dark?). Dark slides get no footer.
    # The graphics-heavy v2 slides (slide_lib_architecture, slide_capability_handoff,
    # slide_one_handler_three_surfaces, slide_lib_registries_v2, slide_agui_event_timeline,
    # slide_production_chain) replace bullet-list versions for stronger visual impact.
    program = [
        (slide_cover,                       "Cover",                  True),
        (slide_agenda,                      "Agenda",                 False),
        # Section 1 — the shift
        (slide_shift,                       "1 · The shift",          False),
        (slide_why_now,                     "1 · The shift",          False),
        (slide_three_protocols,             "1 · The shift",          False),
        # Section 2 — AG-UI
        (slide_agui_background,             "2 · AG-UI",              False),
        (slide_agui_event_timeline,         "2 · AG-UI",              False),  # NEW graphic
        (slide_agui_capabilities,           "2 · AG-UI",              False),
        # Section 3 — Hashbrown
        (slide_hashbrown_background,        "3 · Hashbrown",          False),
        (slide_hashbrown_features,          "3 · Hashbrown",          False),
        (slide_hashbrown_vs_agui,           "3 · Hashbrown",          False),
        (slide_protocol_picker,             "3 · Protocols",          False),  # NEW v4 — when/where
        # Section 4 — A2UI
        (slide_a2ui,                        "4 · A2UI",               False),
        # Section 5 — the library
        (slide_lib_problem,                 "5 · The library",        False),
        (slide_lib_architecture,            "5 · The library",        False),  # rebuilt graphic
        (slide_full_picture,                "5 · The library",        False),  # NEW v4 — full data-flow
        (slide_lib_registries_v2,           "5 · The library",        False),
        (slide_capability_handoff,          "5 · The library",        False),
        (slide_one_handler_three_surfaces,  "5 · The library",        False),
        (slide_lib_backend,                 "5 · The library",        False),
        (slide_lib_mfe,                     "5 · The library",        False),
        (slide_lib_mcp,                     "5 · The library",        False),
        (slide_lib_telemetry,               "5 · The library",        False),
        (slide_lib_schematics,              "5 · The library",        False),
        # Section 6 — examples
        (slide_examples_overview,           "6 · Examples",           False),
        (slide_example_ediscovery,          "6 · Examples",           False),
        (slide_production_chain,            "6 · Examples",           False),
        # Step-reveal sequence — four overlay slides for animated effect
        (slide_run_sequence_step1,          "6 · Examples",           False),  # NEW v4 step-reveal
        (slide_run_sequence_step2,          "6 · Examples",           False),
        (slide_run_sequence_step3,          "6 · Examples",           False),
        (slide_run_sequence_step4,          "6 · Examples",           False),
        # Section 6b — compliance + governance (NEW v5, Phases 5-8)
        (slide_chain_of_custody,            "6 · Compliance",         False),  # NEW v5 P5
        (slide_mcp_reach,                   "7 · MCP reach",          False),  # NEW v5 P6
        (slide_persona_scope_policy,        "8 · Governance",         False),  # NEW v5 P7+8
        (slide_compliance_summary,          "9 · Enterprise readiness", False),  # NEW v5 wrap-up
        # Section 7 — decisions
        (slide_when_to_use_what,            "7 · Decisions",          False),
        # Section 8 — benefits
        (slide_benefits_execs,              "8 · Benefits",           False),
        (slide_benefits_architects,         "8 · Benefits",           False),
        (slide_benefits_devs,               "8 · Benefits",           False),
        # Section 9 — roadmap
        (slide_roadmap,                     "9 · Roadmap",            False),
        # Section 10 — close
        (slide_cta,                         "10 · Call to action",    True),
        (slide_resources,                   "Resources",              False),
    ]

    total = len(program)
    for i, (builder, section, dark) in enumerate(program):
        slide = builder(prs)
        if not dark:
            slide_footer(slide, i + 1, total, section)

    prs.save(OUT)
    print(f"Wrote {OUT}  ({total} slides)")


if __name__ == "__main__":
    main()
